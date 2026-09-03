import streamlit as st
import os
import io
import pandas as pd
from collections import Counter
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter as gcl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ════════════════════════════════════════════════════════════════
#  DEFENSIVEIQ  —  Opponent Offensive Tendency Scouting Report
#  Built for defensive coordinators: upload the opponent's OFFENSIVE
#  playlist (Hudl export) and get a full defensive game-planning
#  package — field-zone tendencies, run/pass concepts, formation
#  tendencies, hash tendencies, down & distance, a call-sheet
#  builder, and a printable game-day call sheet.
# ════════════════════════════════════════════════════════════════

# ── FLEXIBLE COLUMN MAPPING ────────────────────────────────────
# Maps whatever headers a coach's Hudl export uses to the standard
# names the analysis expects. Case / spacing / punctuation tolerant.
COLUMN_ALIASES = {
    "DN":          ["DN", "DOWN", "DWN"],
    "DIST":        ["DIST", "DISTANCE", "DIS", "TO GO", "TOGO"],
    "HASH":        ["HASH", "HASH MARK"],
    "YARD LN":     ["YARD LN", "YARD LINE", "YRDLN", "YARDLINE", "YD LN", "LOS", "FIELD POS"],
    "OFF FORM":    ["OFF FORM", "OFFENSIVE FORMATION", "OFF FORMATION", "FORMATION", "FORM"],
    "CONCEPT":     ["OFF PLAY", "PLAY", "PLAY NAME", "PLAY CALL", "CONCEPT",
                    "PASS CONCEPT", "RUN CONCEPT"],
    "PLAY TYPE":   ["PLAY TYPE", "RUN/PASS", "R/P", "TYPE", "RUNPASS", "ODK"],
    "PLAY DIR":    ["PLAY DIR", "DIRECTION", "DIR", "PLAY DIRECTION"],
    "GN/LS":       ["GN/LS", "GAIN/LOSS", "GAIN", "YARDS", "YDS", "GN LS", "GAINLOSS"],
    "OFF STR":     ["OFF STR", "STRENGTH", "FORM STRENGTH"],
    "BACKFIELD":   ["BACKFIELD", "BACK FIELD", "BACKFLD"],
    "RESULT":      ["RESULT", "RES", "OUTCOME"],
    "QTR":         ["QTR", "QUARTER", "QT", "Q"],
    "PERSONNEL":   ["PERSONNEL", "PERS", "PERSONEL", "GROUPING"],
    "MOTION":      ["MOTION", "MOT", "MOTION DIR"],
    "FORM FAMILY": ["FORM FAMILY", "FORMATION FAMILY", "FAMILY"],
    "FIB":         ["FIB"],
    "PASSER":      ["OPP PASSER", "PASSER", "QB", "QUARTERBACK"],
    "RUSHER":      ["OPP RUSHER", "RUSHER", "BALL CARRIER", "BALLCARRIER", "RB"],
    "RECEIVER":    ["OPP RECEIVER", "RECEIVER", "TARGET", "WR"],
    "BACK DEPTH":  ["BACK DEPTH", "BACKDEPTH", "DEPTH"],
    "OPEN/CLOSE":  ["OPEN/CLOSE", "OPEN/CLOSED", "OPEN CLOSE", "OPENCLOSE"],
}

def _normalize(s):
    return "".join(ch for ch in str(s).upper().strip() if ch.isalnum() or ch == " ").strip()

def map_columns(df):
    """Rename incoming columns to standard names. Returns (df, matched, missing)."""
    norm_incoming = {_normalize(c): c for c in df.columns}
    rename, matched = {}, {}
    for standard, aliases in COLUMN_ALIASES.items():
        for alias in aliases:
            na = _normalize(alias)
            if na in norm_incoming:
                original = norm_incoming[na]
                rename[original] = standard
                matched[standard] = original
                break
    df = df.rename(columns=rename)
    # If two different original columns both map to the same standard name,
    # keep the first and drop the rest so downstream code sees one column.
    df = df.loc[:, ~df.columns.duplicated()]
    missing = [s for s in COLUMN_ALIASES if s not in matched]
    return df, matched, missing

# Columns the analysis genuinely cannot run without
REQUIRED_COLS = ["PLAY TYPE", "YARD LN", "DN", "DIST"]

def check_required(matched):
    return [c for c in REQUIRED_COLS if c not in matched]

st.set_page_config(
    page_title="DefensiveIQ — Opponent Scouting Report",
    page_icon="🛡️",
    layout="wide",
)

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Barlow+Condensed:wght@700;800;900&family=Barlow:wght@400;500;600&family=Share+Tech+Mono&display=swap');
html,body,[class*="css"]{font-family:'Barlow',sans-serif;background-color:#0a1628;color:#f0ede8;}
.stApp{background-color:#0a1628;}
.main-title{font-family:'Barlow Condensed',sans-serif;font-weight:900;font-size:64px;line-height:.95;text-transform:uppercase;color:#f0ede8;margin-bottom:8px;}
.stButton>button{background:#c0392b!important;color:#f0ede8!important;border:none!important;font-family:'Barlow Condensed',sans-serif!important;font-weight:700!important;font-size:16px!important;letter-spacing:.1em!important;text-transform:uppercase!important;padding:12px 32px!important;border-radius:0!important;width:100%!important;}
.stButton>button:hover{background:#a93226!important;}
.stDownloadButton>button{background:#154360!important;color:#f0ede8!important;border:none!important;font-family:'Barlow Condensed',sans-serif!important;font-weight:700!important;font-size:14px!important;letter-spacing:.08em!important;text-transform:uppercase!important;border-radius:0!important;width:100%!important;}
</style>
""", unsafe_allow_html=True)

# ── Helpers ───────────────────────────────────────────────────
def get_zone(y):
    """Standard Hudl field-position convention: negative = own side
    beyond the 50, positive = opponent side."""
    try:
        y = float(y)
        if y <= -1  and y >= -20: return "BZ"
        if y <= -21 and y >= -49: return "OF"
        if y >= 40  and y <= 50:  return "MF"
        if y >= 21  and y <= 39:  return "FZ"
        if y >= 11  and y <= 20:  return "RZ"
        if y >= 1   and y <= 10:  return "GL"
    except Exception:
        pass
    return None

ZONE_LIST  = ["BZ", "OF", "MF", "FZ", "RZ", "GL"]
ZONE_NAMES = {"BZ": "Backed Up  Own 1\u201320", "OF": "Open Field  Own 21\u201349",
              "MF": "Midfield  50\u2013Opp 40", "FZ": "Fringe  Opp 39\u201321",
              "RZ": "Red Zone  Opp 20\u201311", "GL": "Goal Line  Opp 10 and in"}

def top3(plays, key, n=3):
    vals = [str(p.get(key, '')) for p in plays
            if p.get(key) not in (None, '') and str(p.get(key, '')).strip() not in ('', 'nan', 'None')]
    if not vals: return []
    return [{"v": v, "n": c} for v, c in Counter(vals).most_common(n)]

def fmt_top(lst, i):
    return f"{lst[i]['v']} ({lst[i]['n']})" if i < len(lst) else "—"

def pct(n, d): return round(n / d * 100) if d > 0 else 0

def _num(v, default=0.0):
    """Safely parse a number from a spreadsheet cell — handles NaN,
    blanks, and non-numeric junk without poisoning averages."""
    try:
        f = float(v)
        if f != f:  # NaN
            return default
        return f
    except (TypeError, ValueError):
        return default

def is_success(dn, dist, gain, zone, yard_ln, is_td):
    """Success = TD (auto), or gained enough yards for the down.
    Goal line/inside 10: measure yards needed to reach end zone.
    Else: 1st >=40%, 2nd >=50%, 3rd/4th = convert (100%)."""
    if is_td:
        return True
    try:
        dn = int(dn); dist = float(dist); gain = float(gain)
    except Exception:
        return None
    if zone == "GL":
        try:
            to_goal = float(yard_ln)
            if to_goal > 0:
                dist = to_goal
        except Exception:
            pass
    if dist <= 0: dist = 10
    if dn == 1: return gain >= 0.40 * dist
    if dn == 2: return gain >= 0.50 * dist
    if dn >= 3: return gain >= dist
    return None

def load_plays(df):
    """Turn the raw dataframe into a list of play dicts the whole
    app works from. Only Run/Pass snaps with a readable field
    position are kept — special teams / no-plays are dropped."""
    plays = []
    for _, row in df.iterrows():
        pt = str(row.get('PLAY TYPE', '')).strip()
        pt_norm = pt.upper()
        if pt_norm in ('R', 'RUN'): pt = 'Run'
        elif pt_norm in ('P', 'PASS'): pt = 'Pass'
        if pt not in ('Run', 'Pass'):
            continue
        zone = get_zone(row.get('YARD LN', ''))
        if not zone:
            continue
        dn_v   = int(_num(row.get('DN', 0)))
        dist_v = _num(row.get('DIST', 0))
        gain_v = _num(row.get('GN/LS', 0))
        yl_v   = row.get('YARD LN', '')
        result = str(row.get('RESULT', '')).upper()
        is_td  = ('TD' in result) or ('TOUCHDOWN' in result)
        succ   = is_success(dn_v, dist_v, gain_v, zone, yl_v, is_td)
        expl   = (gain_v >= 10) if pt == 'Run' else (gain_v >= 15)
        plays.append({
            'zone':  zone,
            'dn':    dn_v,
            'dist':  dist_v,
            'hash':  str(row.get('HASH', '')).strip().upper()[:1],
            'form':  str(row.get('OFF FORM', '')).strip(),
            'concept': str(row.get('CONCEPT', '')).strip(),
            'dir':   str(row.get('PLAY DIR', '')).strip().upper()[:1],
            'rp':    pt,
            'gnls':  gain_v,
            'strength': str(row.get('OFF STR', '')).strip(),
            'backfield': str(row.get('BACKFIELD', '')).strip(),
            'result': str(row.get('RESULT', '')).strip(),
            'succ':  succ,
            'expl':  expl,
            'td':    is_td,
            'form_family': str(row.get('FORM FAMILY', '')).strip(),
            'fib':   str(row.get('FIB', '')).strip(),
            'passer':   row.get('PASSER', ''),
            'rusher':   row.get('RUSHER', ''),
            'receiver': row.get('RECEIVER', ''),
            'back_depth': str(row.get('BACK DEPTH', '')).strip(),
            'open_close': str(row.get('OPEN/CLOSE', '')).strip(),
        })
    return plays

def _sr(lst):
    v = [p for p in lst if p.get('succ') is not None]
    return round(sum(1 for p in v if p['succ']) / len(v) * 100) if v else None

def _avg(lst):
    return (sum(p['gnls'] for p in lst) / len(lst)) if lst else 0

def dd_bucket(p):
    dn, dist = p['dn'], p['dist']
    if dn == 1: return "1st & 10" if dist >= 8 else "1st & Short"
    if dn == 2:
        if dist >= 7: return "2nd & Long"
        if dist >= 4: return "2nd & Medium"
        return "2nd & Short"
    if dn == 3:
        if dist >= 7: return "3rd & Long"
        if dist >= 4: return "3rd & Medium"
        return "3rd & Short"
    if dn == 4: return "4th Down"
    return "—"

DD_SITS = [
    ("1ST & 10",      lambda p: p['dn'] == 1 and p['dist'] >= 8),
    ("1ST & SHORT",   lambda p: p['dn'] == 1 and p['dist'] < 8),
    ("2ND & LONG",    lambda p: p['dn'] == 2 and p['dist'] >= 7),
    ("2ND & MEDIUM",  lambda p: p['dn'] == 2 and 4 <= p['dist'] <= 6),
    ("2ND & SHORT",   lambda p: p['dn'] == 2 and p['dist'] <= 3),
    ("3RD & LONG",    lambda p: p['dn'] == 3 and p['dist'] >= 7),
    ("3RD & MEDIUM",  lambda p: p['dn'] == 3 and 4 <= p['dist'] <= 6),
    ("3RD & SHORT",   lambda p: p['dn'] == 3 and p['dist'] <= 3),
    ("4TH DOWN",      lambda p: p['dn'] == 4),
    ("RED ZONE",      lambda p: p['zone'] == 'RZ'),
    ("GOAL LINE",     lambda p: p['zone'] == 'GL'),
    ("BACKED UP",     lambda p: p['zone'] == 'BZ'),
]

FZ_SITS = [
    ("1st Down",   lambda p: p['dn'] == 1),
    ("2nd & 7+",   lambda p: p['dn'] == 2 and p['dist'] >= 7),
    ("2nd & 4-6",  lambda p: p['dn'] == 2 and 4 <= p['dist'] <= 6),
    ("2nd & 1-3",  lambda p: p['dn'] == 2 and p['dist'] <= 3),
    ("3rd & 7+",   lambda p: p['dn'] == 3 and p['dist'] >= 7),
    ("3rd & 4-6",  lambda p: p['dn'] == 3 and 4 <= p['dist'] <= 6),
    ("3rd & 1-3",  lambda p: p['dn'] == 3 and p['dist'] <= 3),
    ("4th Down",   lambda p: p['dn'] == 4),
]

def compute_biggest_tendencies(plays, min_n=5, thresh=0.85, top_n=6):
    """Formation/backfield groupings that skew heavily run or pass —
    the headline alerts for a call sheet."""
    out = []
    for key, label in [('form', ''), ('backfield', '')]:
        groups = {}
        for p in plays:
            v = str(p.get(key, '')).strip()
            if v in ('', 'nan', 'None'): continue
            groups.setdefault(v, []).append(p)
        for v, g in groups.items():
            if len(g) < min_n: continue
            n_run = len([p for p in g if p['rp'] == 'Run'])
            run_rate = n_run / len(g)
            if run_rate >= thresh:
                out.append((run_rate * len(g), f"{v} = {round(run_rate*100)}% Run   [{len(g)}]"))
            elif (1 - run_rate) >= thresh:
                out.append(((1 - run_rate) * len(g), f"{v} = {round((1-run_rate)*100)}% Pass   [{len(g)}]"))
    # 3rd & long pass rate
    tl = [p for p in plays if p['dn'] == 3 and p['dist'] >= 7]
    if len(tl) >= 5:
        pr = len([p for p in tl if p['rp'] == 'Pass']) / len(tl)
        out.append((pr * len(tl) * 1.2, f"{round(pr*100)}% Pass on 3rd & Long (7+)   [{len(tl)}]"))
    out.sort(key=lambda t: -t[0])
    return [txt for _, txt in out[:top_n]]

def compute_formation_alerts(plays, min_n=3, top_n=6):
    """Per formation, the single concept that shows up disproportionately."""
    groups = {}
    for p in plays:
        f = p['form']
        if f.strip() in ('', 'nan', 'None'): continue
        groups.setdefault(f, []).append(p)
    alerts = []
    for f, g in groups.items():
        if len(g) < min_n: continue
        cc = Counter(str(p['concept']) for p in g
                     if str(p['concept']).strip() not in ('', 'nan', 'None'))
        if not cc: continue
        concept, n = cc.most_common(1)[0]
        share = n / len(g)
        spread = "spread" if len(g) >= 7 and share < 0.4 else None
        tag = f"{concept} ({round(share*100)}%)" + (f" ({spread})" if spread else "")
        alerts.append((share * len(g), f, tag))
    alerts.sort(key=lambda t: -t[0])
    return [(f, tag) for _, f, tag in alerts[:top_n]]

def compute_heavy_pass_situations(plays):
    heavy = []
    tl = [p for p in plays if p['dn'] == 3 and p['dist'] >= 8]
    if len(tl) >= 3:
        heavy.append(f"3rd & 8+: {pct(len([p for p in tl if p['rp']=='Pass']), len(tl))}% Pass")
    tm = [p for p in plays if p['dn'] == 3 and p['dist'] >= 5]
    if len(tm) >= 3:
        heavy.append(f"2-Min (3rd & 5+): {pct(len([p for p in tm if p['rp']=='Pass']), len(tm))}% Pass")
    sl = [p for p in plays if p['dn'] == 2 and p['dist'] >= 10]
    if len(sl) >= 3:
        heavy.append(f"2nd & 10+: {pct(len([p for p in sl if p['rp']=='Pass']), len(sl))}% Pass")
    return heavy

def _jersey(v):
    """Format a raw passer/rusher/receiver cell as a display label like '#12'."""
    s = str(v).strip()
    if s in ('', 'nan', 'None'): return None
    try:
        f = float(s)
        return f"#{int(f)}"
    except (TypeError, ValueError):
        return f"#{s}"

def compute_player_stats(plays):
    """Build Passing / Rushing / Receiving stat lines per player from
    the RESULT text and the PASSER/RUSHER/RECEIVER columns. Handles
    whatever casing/wording the film uses, and leaves Fum/Int/Drop at
    0 (rendered as '-') when the export doesn't tag those events."""
    passing, rushing, receiving = {}, {}, {}

    for p in plays:
        result_u = p['result'].upper()
        is_incomplete = 'INCOMPLETE' in result_u
        is_complete = ('COMPLETE' in result_u) and not is_incomplete
        is_int = 'INT' in result_u
        is_fum = 'FUM' in result_u
        is_drop = 'DROP' in result_u
        is_td = p['td']

        if p['rp'] == 'Pass':
            passer = _jersey(p.get('passer', ''))
            if passer:
                d = passing.setdefault(passer, {'att': 0, 'cmp': 0, 'yds': 0, 'td': 0,
                                                 'int': 0, 'fum': 0, 'lng': None})
                if is_complete or is_incomplete or is_int:
                    d['att'] += 1
                if is_complete:
                    d['cmp'] += 1
                    d['yds'] += p['gnls']
                    if is_td: d['td'] += 1
                    if d['lng'] is None or p['gnls'] > d['lng']: d['lng'] = p['gnls']
                if is_int: d['int'] += 1
                if is_fum: d['fum'] += 1

            receiver = _jersey(p.get('receiver', ''))
            if receiver:
                d = receiving.setdefault(receiver, {'rec': 0, 'yds': 0, 'td': 0,
                                                      'fum': 0, 'drop': 0, 'lng': None})
                if is_complete:
                    d['rec'] += 1
                    d['yds'] += p['gnls']
                    if is_td: d['td'] += 1
                    if d['lng'] is None or p['gnls'] > d['lng']: d['lng'] = p['gnls']
                if is_drop: d['drop'] += 1
                if is_fum: d['fum'] += 1

        elif p['rp'] == 'Run':
            rusher = _jersey(p.get('rusher', ''))
            if rusher:
                d = rushing.setdefault(rusher, {'att': 0, 'yds': 0, 'td': 0, 'fum': 0, 'lng': None})
                d['att'] += 1
                d['yds'] += p['gnls']
                if d['lng'] is None or p['gnls'] > d['lng']: d['lng'] = p['gnls']
                if is_td: d['td'] += 1
                if is_fum: d['fum'] += 1

    return passing, rushing, receiving

# ── PowerPoint scouting deck (defensive game-plan for a DC) ────
P_NAVY = RGBColor(0x16, 0x21, 0x3E); P_RED = RGBColor(0xC0, 0x39, 0x2B)
P_BLUE = RGBColor(0x1A, 0x52, 0x76); P_TEAL = RGBColor(0x0E, 0x70, 0x60)
P_GOLD = RGBColor(0xC9, 0xA2, 0x27); P_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
P_LGRAY = RGBColor(0xF2, 0xF2, 0xF2); P_DGRAY = RGBColor(0x55, 0x55, 0x55)
P_BLACK = RGBColor(0x11, 0x11, 0x11)
PW, PH = Inches(13.333), Inches(7.5)

def _p_slide(prs, bg=P_WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(1, 0, 0, PW, PH)
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
    sp = r._element; sp.getparent().remove(sp); s.shapes._spTree.insert(2, sp)
    return s

def _p_text(slide, x, y, w, h, text, size=14, color=P_BLACK, bold=False, align=PP_ALIGN.LEFT,
            font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color; run.font.name = font
    return tb

def _p_rect(slide, x, y, w, h, fill, line=None):
    r = slide.shapes.add_shape(1, x, y, w, h)
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if line: r.line.color.rgb = line; r.line.width = Pt(0.75)
    else: r.line.fill.background()
    r.shadow.inherit = False
    return r

def _p_table(slide, x, y, w, rows, col_widths, font_size=11, row_h=Inches(0.34)):
    gt = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, row_h * len(rows)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cw in enumerate(col_widths): gt.columns[ci].width = cw
    for ri, rowvals in enumerate(rows):
        gt.rows[ri].height = row_h
        for ci, (val, opts) in enumerate(rowvals):
            cell = gt.cell(ri, ci)
            cell.margin_left = Pt(4); cell.margin_right = Pt(4)
            cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid(); cell.fill.fore_color.rgb = opts.get('bg', P_WHITE if ri % 2 else P_LGRAY)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = opts.get('align', PP_ALIGN.LEFT)
            run = p.add_run(); run.text = str(val)
            run.font.size = Pt(opts.get('size', font_size))
            run.font.bold = opts.get('bold', False)
            run.font.color.rgb = opts.get('fc', P_BLACK)
            run.font.name = 'Calibri'
    return gt

def build_pptx(plays, opp, week, date, primary_hex="#16213E", accent_hex="#C0392B"):
    def _hex2rgb(h):
        h = h.lstrip('#'); return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    def _lum(h):
        h = h.lstrip('#'); r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return (0.299 * r + 0.587 * g + 0.114 * b) / 255
    PRIMARY = _hex2rgb(primary_hex); ACCENT = _hex2rgb(accent_hex)
    ON_PRIMARY = P_WHITE if _lum(primary_hex) < 0.55 else P_BLACK
    ACCENT_ON_PRIMARY = ACCENT if _lum(accent_hex) > 0.35 else P_WHITE

    prs = Presentation(); prs.slide_width = PW; prs.slide_height = PH
    total = len(plays); opp = opp or "Opponent"
    runs = [p for p in plays if p['rp'] == 'Run']; passes = [p for p in plays if p['rp'] == 'Pass']

    # SLIDE 1 — Title
    s = _p_slide(prs, PRIMARY)
    _p_text(s, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.2),
            "OPPONENT OFFENSE — SCOUTING REPORT", 40, ON_PRIMARY, bold=True, align=PP_ALIGN.CENTER, font="Cambria")
    _p_text(s, Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.9),
            opp.upper(), 32, ACCENT_ON_PRIMARY, bold=True, align=PP_ALIGN.CENTER, font="Cambria")
    sub = "  \u00b7  ".join([x for x in [f"Week {week}" if week else "", date or "", f"{total} plays analyzed"] if x])
    _p_text(s, Inches(0.8), Inches(4.5), Inches(11.7), Inches(0.5), sub, 16,
            RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.CENTER)
    for lx, lbl in [(Inches(1.0), "[ YOUR LOGO ]"), (Inches(9.9), "[ OPP LOGO ]")]:
        _p_rect(s, lx, Inches(0.6), Inches(2.4), Inches(1.3), PRIMARY, line=RGBColor(0x3A, 0x4A, 0x6A))
        _p_text(s, lx, Inches(1.05), Inches(2.4), Inches(0.5), lbl, 11, RGBColor(0x6A, 0x7A, 0x9A),
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _p_text(s, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4),
            "DefensiveIQ  \u00b7  Auto-generated from film — replace logos and decorate freely", 10,
            RGBColor(0x6A, 0x7A, 0x9A), align=PP_ALIGN.CENTER, italic=True)

    # SLIDE 2 — Overview
    s = _p_slide(prs, P_WHITE)
    _p_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
            "OFFENSIVE OVERVIEW", 36, PRIMARY, bold=True, font="Cambria")
    stats = [("PLAYS SCOUTED", str(total), PRIMARY),
             ("RUN %", f"{pct(len(runs), total)}%", P_RED),
             ("PASS %", f"{pct(len(passes), total)}%", P_BLUE),
             ("AVG YDS/PLAY", f"{_avg(plays):.1f}", P_TEAL),
             ("EXPL% (RUN10/PASS15)", f"{round(len([p for p in plays if p['expl']])/total*100) if total else 0}%", P_GOLD)]
    cw = Inches(2.3); gap = Inches(0.15); x0 = Inches(0.6); y0 = Inches(1.5)
    for i, (lbl, val, col) in enumerate(stats):
        x = x0 + (cw + gap) * i
        _p_rect(s, x, y0, cw, Inches(1.8), P_LGRAY)
        _p_text(s, x, y0 + Inches(0.25), cw, Inches(0.9), val, 40, col, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Cambria")
        _p_text(s, x, y0 + Inches(1.25), cw, Inches(0.4), lbl, 11, P_DGRAY, bold=True, align=PP_ALIGN.CENTER)

    _p_text(s, Inches(0.6), Inches(3.7), Inches(12), Inches(0.5), "KEY READS", 20, P_RED, bold=True)
    reads = []
    fc = top3(plays, 'form', 1)
    if fc: reads.append(f"Top formation is {fc[0]['v']} — {fc[0]['n']} snaps")
    cc = top3([p for p in plays if p['rp'] == 'Run'], 'concept', 1)
    if cc: reads.append(f"Top run concept is {cc[0]['v']} — {cc[0]['n']} calls")
    pc = top3([p for p in plays if p['rp'] == 'Pass'], 'concept', 1)
    if pc: reads.append(f"Top pass concept is {pc[0]['v']} — {pc[0]['n']} calls")
    for txt in compute_biggest_tendencies(plays, top_n=2):
        reads.append(txt)
    yy = Inches(4.35)
    for txt in reads[:6]:
        _p_rect(s, Inches(0.7), yy + Inches(0.05), Inches(0.22), Inches(0.22), P_RED)
        _p_text(s, Inches(1.05), yy, Inches(11), Inches(0.4), txt, 16, P_BLACK, bold=True)
        yy += Inches(0.5)

    # SLIDE 5 — Formation Tendencies
    s = _p_slide(prs, P_WHITE)
    _p_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
            "FORMATION TENDENCIES", 34, RGBColor(0x4A, 0x23, 0x5A), bold=True, font="Cambria")
    fgroups = {}
    for p in plays:
        f = p['form']
        if f.strip() in ('', 'nan', 'None'): continue
        fgroups.setdefault(f, []).append(p)
    ranked = sorted(fgroups.items(), key=lambda kv: -len(kv[1]))
    ranked = [(f, g) for f, g in ranked if len(g) >= 3][:9]
    rows = [[("FORMATION", {'bg': RGBColor(0x4A, 0x23, 0x5A), 'fc': P_WHITE, 'bold': True}),
             ("SNAPS", {'bg': RGBColor(0x4A, 0x23, 0x5A), 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
             ("RUN%", {'bg': RGBColor(0x4A, 0x23, 0x5A), 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
             ("TOP RUN CONCEPT", {'bg': RGBColor(0x4A, 0x23, 0x5A), 'fc': P_WHITE, 'bold': True}),
             ("TOP PASS CONCEPT", {'bg': RGBColor(0x4A, 0x23, 0x5A), 'fc': P_WHITE, 'bold': True})]]
    for i, (f, g) in enumerate(ranked):
        bg = P_LGRAY if i % 2 == 0 else P_WHITE
        gr = [p for p in g if p['rp'] == 'Run']; gp = [p for p in g if p['rp'] == 'Pass']
        rows.append([(f, {'bg': bg, 'bold': True, 'size': 10}),
                     (str(len(g)), {'bg': bg, 'align': PP_ALIGN.CENTER}),
                     (f"{pct(len(gr), len(g))}%", {'bg': bg, 'fc': P_RED, 'bold': True, 'align': PP_ALIGN.CENTER}),
                     (fmt_top(top3(gr, 'concept', 1), 0), {'bg': bg, 'size': 10}),
                     (fmt_top(top3(gp, 'concept', 1), 0), {'bg': bg, 'size': 10})])
    if len(rows) == 1:
        _p_text(s, Inches(0.6), Inches(2), Inches(12), Inches(0.5),
                "Not enough formation data tagged.", 14, P_DGRAY, italic=True)
    else:
        _p_table(s, Inches(0.6), Inches(1.5), Inches(12.1), rows,
                 [Inches(3.2), Inches(1.5), Inches(1.5), Inches(2.95), Inches(2.95)], row_h=Inches(0.5))

    # SLIDE 6 — Down & Distance keys
    s = _p_slide(prs, P_WHITE)
    _p_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
            "DOWN & DISTANCE KEYS", 34, P_TEAL, bold=True, font="Cambria")
    rows = [[("SITUATION", {'bg': P_TEAL, 'fc': P_WHITE, 'bold': True}),
             ("SNAPS", {'bg': P_TEAL, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
             ("RUN%", {'bg': P_TEAL, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
             ("PASS%", {'bg': P_TEAL, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
             ("TOP RUN", {'bg': P_TEAL, 'fc': P_WHITE, 'bold': True}),
             ("TOP PASS", {'bg': P_TEAL, 'fc': P_WHITE, 'bold': True})]]
    for i, (lbl, fn) in enumerate(DD_SITS):
        sp = [p for p in plays if fn(p)]
        if not sp: continue
        bg = P_LGRAY if i % 2 == 0 else P_WHITE
        sr = [p for p in sp if p['rp'] == 'Run']; spa = [p for p in sp if p['rp'] == 'Pass']
        rows.append([(lbl, {'bg': bg, 'bold': True, 'size': 10}),
                     (str(len(sp)), {'bg': bg, 'align': PP_ALIGN.CENTER}),
                     (f"{pct(len(sr), len(sp))}%", {'bg': bg, 'fc': P_RED, 'bold': True, 'align': PP_ALIGN.CENTER}),
                     (f"{pct(len(spa), len(sp))}%", {'bg': bg, 'fc': P_BLUE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                     (fmt_top(top3(sr, 'concept', 1), 0), {'bg': bg, 'size': 10}),
                     (fmt_top(top3(spa, 'concept', 1), 0), {'bg': bg, 'size': 10})])
    _p_table(s, Inches(0.6), Inches(1.5), Inches(12.1), rows,
             [Inches(2.7), Inches(1.5), Inches(1.5), Inches(1.5), Inches(2.45), Inches(2.45)], row_h=Inches(0.47))

    # ── Helper: concept-level detail for the Top Concepts slide ──
    def _concept_detail(rp_filter, n=5):
        grp = {}
        for p in plays:
            if p['rp'] != rp_filter: continue
            c = str(p['concept']).strip()
            if c in ('', 'nan', 'None'): continue
            grp.setdefault(c, []).append(p)
        ranked = sorted(grp.items(), key=lambda kv: -len(kv[1]))[:n]
        out = []
        for concept, g in ranked:
            expl_pct = len([p for p in g if p['expl']]) / len(g) if g else 0
            sr = _sr(g)
            tf = top3(g, 'form', 1)
            out.append((concept, len(g), _avg(g), expl_pct, sr, tf[0]['v'] if tf else "—"))
        return out

    # ── Helper: group breakdown (Form Family / FIB) for its slide ──
    def _group_detail(key, empty_label, n=None):
        groups = {}
        for p in plays:
            v = str(p.get(key, '')).strip()
            if v in ('', 'nan', 'None'): v = empty_label
            groups.setdefault(v, []).append(p)
        ranked = sorted(groups.items(), key=lambda kv: -len(kv[1]))
        total_groups = len(ranked)
        if n: ranked = ranked[:n]
        out = []
        for v, g in ranked:
            gr = [p for p in g if p['rp'] == 'Run']; gp = [p for p in g if p['rp'] == 'Pass']
            out.append((v, len(g), len(gr) / len(g) if g else 0, len(gp) / len(g) if g else 0,
                        fmt_top(top3(gr, 'concept', 1), 0), fmt_top(top3(gp, 'concept', 1), 0),
                        fmt_top(top3(g, 'form', 1), 0)))
        return out, total_groups

    # SLIDE 7 — Top Run & Pass Concepts (detailed)
    s = _p_slide(prs, P_WHITE)
    _p_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
            "TOP RUN & PASS CONCEPTS", 32, RGBColor(0x4A, 0x23, 0x5A), bold=True, font="Cambria")

    def _concept_table(y, title_txt, rp_filter, color):
        _p_text(s, Inches(0.6), y, Inches(12), Inches(0.35), title_txt, 16, color, bold=True)
        data = _concept_detail(rp_filter, 5)
        rows = [[("CONCEPT", {'bg': color, 'fc': P_WHITE, 'bold': True}),
                 ("CALLED", {'bg': color, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("AVG YD", {'bg': color, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("EXPL%", {'bg': color, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("SUCC%", {'bg': color, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("TOP FORM", {'bg': color, 'fc': P_WHITE, 'bold': True})]]
        for i, (concept, called, avg, expl, succ, form) in enumerate(data):
            bg = P_LGRAY if i % 2 == 0 else P_WHITE
            rows.append([(concept, {'bg': bg, 'bold': True, 'size': 11}),
                         (str(called), {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (f"{avg:.1f}", {'bg': bg, 'fc': P_TEAL, 'bold': True, 'align': PP_ALIGN.CENTER}),
                         (f"{round(expl*100)}%", {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (f"{succ}%" if succ is not None else "—", {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (form, {'bg': bg, 'size': 11})])
        if len(rows) == 1:
            _p_text(s, Inches(0.6), y + Inches(0.4), Inches(12), Inches(0.4),
                    "Not enough tagged data.", 13, P_DGRAY, italic=True)
            return
        _p_table(s, Inches(0.6), y + Inches(0.4), Inches(12.1), rows,
                 [Inches(3.2), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6), Inches(2.5)],
                 row_h=Inches(0.4))

    _concept_table(Inches(1.1), "TOP 5 RUN CONCEPTS", "Run", P_RED)
    _concept_table(Inches(4.15), "TOP 5 PASS CONCEPTS", "Pass", P_BLUE)

    # SLIDE 8 — Top 5 Biggest Tendencies
    s = _p_slide(prs, PRIMARY)
    _p_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
            "TOP 5 BIGGEST TENDENCIES", 32, ACCENT_ON_PRIMARY, bold=True, font="Cambria")
    big5 = compute_biggest_tendencies(plays, top_n=5)
    yy = Inches(1.6)
    for i, txt in enumerate(big5):
        _p_rect(s, Inches(0.8), yy, Inches(0.5), Inches(0.55), ACCENT)
        _p_text(s, Inches(0.8), yy, Inches(0.5), Inches(0.55), str(i + 1), 20, P_WHITE,
                bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _p_text(s, Inches(1.45), yy, Inches(10.8), Inches(0.55), txt, 17, P_WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        yy += Inches(0.85)
    if not big5:
        _p_text(s, Inches(0.8), Inches(2), Inches(11), Inches(0.5),
                "Not enough tagged snaps to compute tendencies yet.", 14, RGBColor(0x8A, 0x9A, 0xBA), italic=True)
    _p_text(s, Inches(0.6), Inches(6.95), Inches(12), Inches(0.4),
            "Tendencies require a minimum sample size — verify against film before installing.", 10,
            RGBColor(0x8A, 0x9A, 0xBA), italic=True)

    # SLIDE 9 — Form Family & FIB Tendencies
    s = _p_slide(prs, P_WHITE)
    _p_text(s, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
            "FORM FAMILY & FIB TENDENCIES", 30, RGBColor(0x6C, 0x34, 0x83), bold=True, font="Cambria")

    def _group_table(y, title_txt, key, empty_label, color, n=None):
        _p_text(s, Inches(0.6), y, Inches(12), Inches(0.35), title_txt, 16, color, bold=True)
        data, total_groups = _group_detail(key, empty_label, n)
        rows = [[("GROUP", {'bg': color, 'fc': P_WHITE, 'bold': True}),
                 ("SNAPS", {'bg': color, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("RUN%", {'bg': color, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("PASS%", {'bg': color, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("TOP RUN", {'bg': color, 'fc': P_WHITE, 'bold': True}),
                 ("TOP PASS", {'bg': color, 'fc': P_WHITE, 'bold': True}),
                 ("TOP FORM", {'bg': color, 'fc': P_WHITE, 'bold': True})]]
        for i, (v, snaps, runp, passp, toprun, toppass, topform) in enumerate(data):
            bg = P_LGRAY if i % 2 == 0 else P_WHITE
            rows.append([(v, {'bg': bg, 'bold': True, 'size': 10}),
                         (str(snaps), {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (f"{round(runp*100)}%", {'bg': bg, 'fc': P_RED, 'bold': True, 'align': PP_ALIGN.CENTER}),
                         (f"{round(passp*100)}%", {'bg': bg, 'fc': P_BLUE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                         (toprun, {'bg': bg, 'size': 9}),
                         (toppass, {'bg': bg, 'size': 9}),
                         (topform, {'bg': bg, 'size': 9})])
        if len(rows) == 1:
            _p_text(s, Inches(0.6), y + Inches(0.4), Inches(12), Inches(0.4),
                    "Not enough tagged data.", 13, P_DGRAY, italic=True)
            return
        _p_table(s, Inches(0.6), y + Inches(0.4), Inches(12.1), rows,
                 [Inches(2.0), Inches(1.2), Inches(1.2), Inches(1.2), Inches(2.15), Inches(2.15), Inches(2.2)],
                 row_h=Inches(0.38))
        if n and total_groups > n:
            _p_text(s, Inches(0.6), y + Inches(0.4) + Inches(0.38) * len(rows), Inches(12), Inches(0.3),
                    f"+ {total_groups - n} more — see the Form Family / FIB tabs in the workbook.", 10, P_DGRAY, italic=True)

    _group_table(Inches(1.1), "FORM FAMILY", 'form_family', "(Blank)", RGBColor(0x6C, 0x34, 0x83), n=6)
    _group_table(Inches(4.7), "FIB", 'fib', "Not FIB", RGBColor(0x78, 0x42, 0x12))

    # SLIDE 10 — Stats (Passing / Rushing / Receiving)
    s = _p_slide(prs, P_WHITE)
    _p_text(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.6),
            "PLAYER STATS", 32, PRIMARY, bold=True, font="Cambria")
    passing, rushing, receiving = compute_player_stats(plays)

    def _dash(val):
        return str(val) if val else "—"

    yy = Inches(1.0)
    _p_text(s, Inches(0.6), yy, Inches(12), Inches(0.3), "PASSING", 14, P_RED, bold=True)
    yy += Inches(0.3)
    pass_all = sorted(passing.items(), key=lambda kv: -kv[1]['yds'])
    pass_ranked = pass_all[:3]
    if pass_ranked:
        rows = [[("PLAYER", {'bg': P_RED, 'fc': P_WHITE, 'bold': True}),
                 ("CMP", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("ATT", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("%", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("YDS", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("TD", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("INT", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("RAT", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER})]]
        for i, (player, d) in enumerate(pass_ranked):
            bg = P_LGRAY if i % 2 == 0 else P_WHITE
            att, cmp_, yds, td, intc = d['att'], d['cmp'], d['yds'], d['td'], d['int']
            rating = (8.4 * yds + 330 * td + 100 * cmp_ - 200 * intc) / att if att else 0
            rows.append([(player, {'bg': bg, 'bold': True}),
                         (str(cmp_), {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (str(att), {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (f"{round(cmp_/att*100) if att else 0}%", {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (str(yds), {'bg': bg, 'fc': P_TEAL, 'bold': True, 'align': PP_ALIGN.CENTER}),
                         (_dash(td), {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (_dash(intc), {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (f"{rating:.1f}", {'bg': bg, 'fc': P_RED, 'bold': True, 'align': PP_ALIGN.CENTER})])
        _p_table(s, Inches(0.6), yy, Inches(12.1), rows,
                 [Inches(2.4), Inches(1.2), Inches(1.2), Inches(1.4), Inches(1.5), Inches(1.2), Inches(1.2), Inches(2.0)],
                 row_h=Inches(0.3))
        yy += Inches(0.3) * len(rows)
        if len(pass_all) > len(pass_ranked):
            _p_text(s, Inches(0.6), yy, Inches(12), Inches(0.22),
                    f"+ {len(pass_all) - len(pass_ranked)} more — see the Stats tab in the workbook.", 9, P_DGRAY, italic=True)
        yy += Inches(0.22) + Inches(0.15)
    else:
        _p_text(s, Inches(0.6), yy, Inches(12), Inches(0.3), "No passer data tagged.", 12, P_DGRAY, italic=True)
        yy += Inches(0.37)

    _p_text(s, Inches(0.6), yy, Inches(12), Inches(0.3), "RUSHING", 14, P_RED, bold=True)
    yy += Inches(0.3)
    rush_all = sorted(rushing.items(), key=lambda kv: -kv[1]['yds'])
    rush_ranked = rush_all[:4]
    if rush_ranked:
        rows = [[("PLAYER", {'bg': P_RED, 'fc': P_WHITE, 'bold': True}),
                 ("ATT", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("YDS", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("AVG", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("LNG", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("TD", {'bg': P_RED, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER})]]
        for i, (player, d) in enumerate(rush_ranked):
            bg = P_LGRAY if i % 2 == 0 else P_WHITE
            att, yds, td, lng = d['att'], d['yds'], d['td'], d['lng']
            rows.append([(player, {'bg': bg, 'bold': True}),
                         (str(att), {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (str(yds), {'bg': bg, 'fc': P_TEAL, 'bold': True, 'align': PP_ALIGN.CENTER}),
                         (f"{yds/att:.1f}" if att else "—", {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (str(lng) if lng is not None else "—", {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (_dash(td), {'bg': bg, 'align': PP_ALIGN.CENTER})])
        _p_table(s, Inches(0.6), yy, Inches(12.1), rows,
                 [Inches(3.0), Inches(1.6), Inches(1.9), Inches(1.9), Inches(1.7), Inches(2.0)],
                 row_h=Inches(0.3))
        yy += Inches(0.3) * len(rows)
        if len(rush_all) > len(rush_ranked):
            _p_text(s, Inches(0.6), yy, Inches(12), Inches(0.22),
                    f"+ {len(rush_all) - len(rush_ranked)} more — see the Stats tab in the workbook.", 9, P_DGRAY, italic=True)
        yy += Inches(0.22) + Inches(0.15)
    else:
        _p_text(s, Inches(0.6), yy, Inches(12), Inches(0.3), "No rusher data tagged.", 12, P_DGRAY, italic=True)
        yy += Inches(0.37)

    _p_text(s, Inches(0.6), yy, Inches(12), Inches(0.3), "RECEIVING", 14, P_BLUE, bold=True)
    yy += Inches(0.3)
    rec_all = sorted(receiving.items(), key=lambda kv: -kv[1]['yds'])
    rec_ranked = rec_all[:5]
    if rec_ranked:
        rows = [[("PLAYER", {'bg': P_BLUE, 'fc': P_WHITE, 'bold': True}),
                 ("REC", {'bg': P_BLUE, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("YDS", {'bg': P_BLUE, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("AVG", {'bg': P_BLUE, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("LNG", {'bg': P_BLUE, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER}),
                 ("TD", {'bg': P_BLUE, 'fc': P_WHITE, 'bold': True, 'align': PP_ALIGN.CENTER})]]
        for i, (player, d) in enumerate(rec_ranked):
            bg = P_LGRAY if i % 2 == 0 else P_WHITE
            rec, yds, td, lng = d['rec'], d['yds'], d['td'], d['lng']
            rows.append([(player, {'bg': bg, 'bold': True}),
                         (str(rec), {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (str(yds), {'bg': bg, 'fc': P_TEAL, 'bold': True, 'align': PP_ALIGN.CENTER}),
                         (f"{yds/rec:.1f}" if rec else "—", {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (str(lng) if lng is not None else "—", {'bg': bg, 'align': PP_ALIGN.CENTER}),
                         (_dash(td), {'bg': bg, 'align': PP_ALIGN.CENTER})])
        _p_table(s, Inches(0.6), yy, Inches(12.1), rows,
                 [Inches(3.0), Inches(1.6), Inches(1.9), Inches(1.9), Inches(1.7), Inches(2.0)],
                 row_h=Inches(0.28))
        yy += Inches(0.28) * len(rows)
        if len(rec_all) > len(rec_ranked):
            _p_text(s, Inches(0.6), yy, Inches(12), Inches(0.22),
                    f"+ {len(rec_all) - len(rec_ranked)} more — see the Stats tab in the workbook.", 9, P_DGRAY, italic=True)
    else:
        _p_text(s, Inches(0.6), yy, Inches(12), Inches(0.3), "No receiver data tagged.", 12, P_DGRAY, italic=True)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    return buf.getvalue()

# ── Excel Builder ─────────────────────────────────────────────
def build_excel(plays, opp, week, date):
    FN = "Arial"; CW = "FFFFFFFF"; CL = "FFF5F5F5"; CB = "FF16213E"
    CBl = "FF1A5276"; CTe = "FF0E7060"; CPu = "FF4A235A"; CR = "FFC0392B"
    CRB = "FFFDE8E8"; CPB = "FFE8F0FE"
    CDG = "FF555555"; CGr = "FF1E8449"; CYB = "FFFFFBE6"
    ZONE_BG  = {"BZ": "FFFDE8E8", "OF": "FFE8F0FE", "MF": "FFE8F8E8",
                "FZ": "FFFFFBE6", "RZ": "FFFCE4EC", "GL": "FFEDE7F6"}
    ZONE_HDR = {"BZ": CR, "OF": CBl, "MF": CTe, "FZ": "FF7D6608", "RZ": CR, "GL": CPu}

    def fil(c): return PatternFill("solid", fgColor=c)
    def bdr():
        s = Side(style="thin", color="FFB0B0B0")
        return Border(left=s, right=s, top=s, bottom=s)
    def sc(ws, r, c, val="", bold=False, sz=10, fc=CB, bg=None, h="center", v="center", wrap=False, fmt=None):
        cell = ws.cell(row=r, column=c, value=val)
        cell.font = Font(name=FN, bold=bold, size=sz, color=fc)
        if bg: cell.fill = fil(bg)
        cell.alignment = Alignment(horizontal=h, vertical=v, wrap_text=wrap)
        cell.border = bdr()
        if fmt: cell.number_format = fmt
        return cell
    def hdr(ws, r, c, txt, bg=CBl, fc=CW, sz=9, wrap=True, span=1):
        cell = sc(ws, r, c, txt, bold=True, sz=sz, fc=fc, bg=bg, wrap=wrap)
        if span > 1: ws.merge_cells(start_row=r, start_column=c, end_row=r, end_column=c + span - 1)
        return cell
    def banner(ws, r, txt, nc, bg=CB, fc=CW, sz=13, ht=30):
        ws.merge_cells(start_row=r, start_column=1, end_row=r, end_column=nc)
        c = ws.cell(row=r, column=1, value=txt)
        c.font = Font(name=FN, bold=True, size=sz, color=fc)
        c.fill = fil(bg); c.alignment = Alignment(horizontal="center", vertical="center")
        ws.row_dimensions[r].height = ht
    def widths(ws, lst):
        for i, w in enumerate(lst, 1): ws.column_dimensions[gcl(i)].width = w

    total = len(plays)
    runs = [p for p in plays if p['rp'] == 'Run']
    passes = [p for p in plays if p['rp'] == 'Pass']

    def top3_str(lst, key, n=3):
        vals = [str(p.get(key, '')) for p in lst
                if p.get(key) not in (None, '') and str(p.get(key, '')).strip() not in ('', 'nan', 'None')]
        if not vals: return ["—"] * n
        counts = Counter(vals).most_common(n)
        res = [f"{v} ({c})" for v, c in counts]
        while len(res) < n: res.append("—")
        return res

    wb2 = Workbook()

    # ── Tab 1: Film Log ──────────────────────────────────────
    ws_log = wb2.active; ws_log.title = "1. Film Log"
    ws_log.sheet_properties.tabColor = "C0392B"
    ws_log.sheet_view.showGridLines = False
    log_cols = [('QTR', 6), ('DN', 6), ('DIST', 6), ('HASH', 6), ('YARD LN', 9), ('ZONE', 10),
                ('OFF FORM', 20), ('OFF STR', 8), ('BACKFIELD', 12), ('PLAY DIR', 9),
                ('PLAY TYPE', 10), ('CONCEPT', 20), ('GN/LS', 8), ('RESULT', 12)]
    widths(ws_log, [w for _, w in log_cols])
    banner(ws_log, 1, "DEFENSIVEIQ — FILM LOG  |  CONCEPT = opponent's play/concept name for ALL plays",
           len(log_cols), bg=CR, sz=10, ht=26)
    ws_log.row_dimensions[2].height = 20
    for ci, (col, _) in enumerate(log_cols, 1):
        bg = "FF000088" if col == "ZONE" else (CR if col in ("PLAY TYPE", "CONCEPT") else CBl)
        hdr(ws_log, 2, ci, col, bg=bg, sz=8)
    for ri, p in enumerate(plays):
        r = ri + 3; ws_log.row_dimensions[r].height = 15
        bg = CL if ri % 2 == 0 else CW
        vals = {'QTR': '', 'DN': p['dn'], 'DIST': p['dist'], 'HASH': p['hash'], 'YARD LN': '',
                'ZONE': p['zone'], 'OFF FORM': p['form'], 'OFF STR': p['strength'],
                'BACKFIELD': p['backfield'], 'PLAY DIR': p['dir'], 'PLAY TYPE': p['rp'],
                'CONCEPT': p['concept'], 'GN/LS': p['gnls'], 'RESULT': p['result']}
        for ci, (col, _) in enumerate(log_cols, 1):
            zbg = "FFE8F4FD" if col == "ZONE" else ("FFFDE8E8" if col in ("PLAY TYPE", "CONCEPT") else bg)
            sc(ws_log, r, ci, vals.get(col, ''), sz=9, bg=zbg, fc="FF000000",
               h="left" if col in ("OFF FORM", "CONCEPT", "RESULT", "BACKFIELD") else "center")
    ws_log.freeze_panes = "A3"

    # ── Tab 2: Field Zone Tendencies ─────────────────────────
    ws2 = wb2.create_sheet("2. Field Zone Tendencies")
    ws2.sheet_properties.tabColor = "1A5276"; ws2.sheet_view.showGridLines = False
    NC2 = 9
    widths(ws2, [10, 22] + [10] * (NC2 - 2))
    banner(ws2, 1, "FIELD ZONE TENDENCIES", NC2, bg=CB, sz=13, ht=28)
    ws2.merge_cells(f"A2:{gcl(NC2)}2")
    leg = ws2.cell(row=2, column=1, value="  Gray = plays   Red% = Run   Blue% = Pass   Yellow = Call Idea")
    leg.font = Font(name=FN, size=8, italic=True, color=CDG); leg.fill = fil("FFF0F0F0")
    leg.alignment = Alignment(horizontal="left", vertical="center")
    ws2.row_dimensions[2].height = 14
    row = 3
    hdr(ws2, row, 1, "FIELD ZONE", bg=CBl, sz=9); hdr(ws2, row, 2, "METRIC", bg=CBl, sz=9)
    for ci, (lbl, _fn) in enumerate(FZ_SITS, 3):
        hdr(ws2, row, ci, lbl, bg=CBl, sz=9)
    row += 1
    for zcode in ZONE_LIST:
        zp = [p for p in plays if p['zone'] == zcode]
        ws2.merge_cells(start_row=row, start_column=1, end_row=row, end_column=NC2)
        c = ws2.cell(row=row, column=1, value=f"  {zcode}  \u00b7  {ZONE_NAMES[zcode]}  ({len(zp)} plays)")
        c.font = Font(name=FN, bold=True, size=11, color=CW); c.fill = fil(ZONE_HDR[zcode])
        c.alignment = Alignment(horizontal="left", vertical="center")
        ws2.row_dimensions[row].height = 20
        row += 1
        metric_rows = ["Plays", "Run %", "Pass %", "Runs", "Passes", "\u25b6 Call Idea"]
        vals_by_sit = []
        for _lbl, fn in FZ_SITS:
            sp = [p for p in zp if fn(p)]
            sr = [p for p in sp if p['rp'] == 'Run']; spa = [p for p in sp if p['rp'] == 'Pass']
            vals_by_sit.append((len(sp), len(sr), len(spa)))
        for mi, mlbl in enumerate(metric_rows):
            bg = CL if mi % 2 == 0 else CW
            sc(ws2, row, 1, "", bg=bg)
            sc(ws2, row, 2, mlbl, bold=True, sz=9, fc=CB, bg=bg, h="left")
            for ci, (n, nr, npz) in enumerate(vals_by_sit, 3):
                if mlbl == "Plays": v = n if n else ""
                elif mlbl == "Run %": v = round(nr / n, 2) if n else ""
                elif mlbl == "Pass %": v = round(npz / n, 2) if n else ""
                elif mlbl == "Runs": v = nr if n else ""
                elif mlbl == "Passes": v = npz if n else ""
                else: v = ""
                cbg = CYB if mlbl == "\u25b6 Call Idea" else bg
                fmt = "0%" if "%" in mlbl else ("0" if v != "" else "General")
                fc = "FFC0392B" if mlbl == "Run %" else ("FF00008B" if mlbl == "Pass %" else "FF000000")
                sc(ws2, row, ci, v, sz=9, fc=fc, bg=cbg, fmt=fmt)
            row += 1
        row += 1
    ws2.freeze_panes = "C4"

    # ── Tab 3: Run Tendencies (by zone) ──────────────────────
    ws3 = wb2.create_sheet("3. Run Tendencies")
    ws3.sheet_properties.tabColor = "C0392B"; ws3.sheet_view.showGridLines = False
    NC3 = 13
    widths(ws3, [8, 9, 8, 10, 18, 18, 18, 16, 16, 16, 8, 8, 8])
    banner(ws3, 1, "RUN TENDENCIES  —  Auto-calculated from film", NC3, bg=CR, sz=13, ht=28)
    for c, txt, bg, span in [(1, "ZONE", CB, 1), (2, "COUNTS", CB, 2), (4, "TOP FORMATIONS", CR, 3),
                              (7, "TOP RUN CONCEPTS", CR, 3), (10, "TOP DIRECTION", CR, 3)]:
        hdr(ws3, 2, c, txt, bg=bg, sz=8, span=span)
    for c, txt, bg in [(1, "Zone", CB), (2, "Run\nCount", CB), (3, "Run %", CB), (4, "Expl.\nRuns", CB),
                       (5, "#1 Formation", CR), (6, "#2 Formation", CR), (7, "#3 Formation", CR),
                       (8, "#1 Concept", CR), (9, "#2 Concept", CR), (10, "#3 Concept", CR),
                       (11, "#1 Dir", CR), (12, "#2 Dir", CR), (13, "#3 Dir", CR)]:
        hdr(ws3, 3, c, txt, bg=bg, sz=8, wrap=True)
    for ri, zcode in enumerate(ZONE_LIST):
        r = ri + 4; zbg = ZONE_BG[zcode]; zhdr = ZONE_HDR[zcode]
        zr = [p for p in plays if p['zone'] == zcode and p['rp'] == 'Run']
        ws3.row_dimensions[r].height = 26
        sc(ws3, r, 1, zcode, bold=True, sz=11, fc=CW, bg=zhdr)
        sc(ws3, r, 2, len(zr), bold=True, sz=11, fc="FF000000", bg=zbg, fmt="0")
        allz = [p for p in plays if p['zone'] == zcode]
        sc(ws3, r, 3, round(len(zr) / len(allz), 2) if allz else "", bold=True, sz=10, fc="FF8B0000", bg=CRB, fmt="0%")
        sc(ws3, r, 4, len([p for p in zr if p['expl']]), sz=10, fc="FF0E7060", bg="FFE8F8E8", fmt="0")
        t3f = top3_str(zr, 'form'); t3c = top3_str(zr, 'concept'); t3d = top3_str(zr, 'dir', 2) + ["—"]
        for i, cn in enumerate([5, 6, 7]): sc(ws3, r, cn, t3f[i], sz=9, bg=zbg, wrap=True)
        for i, cn in enumerate([8, 9, 10]): sc(ws3, r, cn, t3c[i], sz=9, bg=zbg, wrap=True)
        for i, cn in enumerate([11, 12, 13]): sc(ws3, r, cn, t3d[i], sz=9, bg=zbg, wrap=True)
    ws3.freeze_panes = "B4"

    # ── Tab 4: Pass Tendencies (by zone) ─────────────────────
    ws4 = wb2.create_sheet("4. Pass Tendencies")
    ws4.sheet_properties.tabColor = "1A5276"; ws4.sheet_view.showGridLines = False
    NC4 = 13
    widths(ws4, [8, 9, 8, 10, 18, 18, 18, 16, 16, 16, 12, 12, 12])
    banner(ws4, 1, "PASS TENDENCIES  —  Auto-calculated from film", NC4, bg=CBl, sz=13, ht=28)
    for c, txt, bg, span in [(1, "ZONE", CB, 1), (2, "COUNTS", CB, 2), (4, "TOP FORMATIONS ON PASS", CBl, 3),
                              (7, "TOP PASS CONCEPTS", CBl, 3), (10, "TOP BACKFIELD", CBl, 3)]:
        hdr(ws4, 2, c, txt, bg=bg, sz=8, span=span)
    for c, txt, bg in [(1, "Zone", CB), (2, "Pass\nCount", CB), (3, "Pass %", CB), (4, "Expl.\nPasses", CB),
                       (5, "#1 Formation", CBl), (6, "#2 Formation", CBl), (7, "#3 Formation", CBl),
                       (8, "#1 Concept", CBl), (9, "#2 Concept", CBl), (10, "#3 Concept", CBl),
                       (11, "#1 Backfield", CBl), (12, "#2 Backfield", CBl), (13, "#3 Backfield", CBl)]:
        hdr(ws4, 3, c, txt, bg=bg, sz=8, wrap=True)
    for ri, zcode in enumerate(ZONE_LIST):
        r = ri + 4; zbg = ZONE_BG[zcode]; zhdr = ZONE_HDR[zcode]
        zp = [p for p in plays if p['zone'] == zcode and p['rp'] == 'Pass']
        allz = [p for p in plays if p['zone'] == zcode]
        ws4.row_dimensions[r].height = 26
        sc(ws4, r, 1, zcode, bold=True, sz=11, fc=CW, bg=zhdr)
        sc(ws4, r, 2, len(zp), bold=True, sz=11, fc="FF000000", bg=zbg, fmt="0")
        sc(ws4, r, 3, round(len(zp) / len(allz), 2) if allz else "", bold=True, sz=10, fc="FF00008B", bg=CPB, fmt="0%")
        sc(ws4, r, 4, len([p for p in zp if p['expl']]), sz=10, fc="FF0E7060", bg="FFE8F8E8", fmt="0")
        t3f = top3_str(zp, 'form'); t3c = top3_str(zp, 'concept'); t3b = top3_str(zp, 'backfield')
        for i, cn in enumerate([5, 6, 7]): sc(ws4, r, cn, t3f[i], sz=9, bg=zbg, wrap=True)
        for i, cn in enumerate([8, 9, 10]): sc(ws4, r, cn, t3c[i], sz=9, bg=zbg, wrap=True)
        for i, cn in enumerate([11, 12, 13]): sc(ws4, r, cn, t3b[i], sz=9, bg=zbg, wrap=True)
    ws4.freeze_panes = "B4"

    # ── Tab 5: Hash Tendencies ───────────────────────────────
    ws5 = wb2.create_sheet("5. Hash Tendencies")
    ws5.sheet_properties.tabColor = "6C3483"; ws5.sheet_view.showGridLines = False
    NC5 = 13
    widths(ws5, [22, 9, 9, 9, 9, 9, 9, 9, 9, 9, 18, 18, 18])
    banner(ws5, 1, "HASH TENDENCIES  —  Left Hash \u00b7 Middle \u00b7 Right Hash", NC5, bg=CPu, sz=13, ht=28)
    for c, txt, bg in [(1, "FIELD ZONE", CB), (2, "L Plays", "FF6C3483"), (3, "L Run%", "FF6C3483"), (4, "L Pass%", "FF6C3483"),
                       (5, "M Plays", CBl), (6, "M Run%", CBl), (7, "M Pass%", CBl),
                       (8, "R Plays", "FF784212"), (9, "R Run%", "FF784212"), (10, "R Pass%", "FF784212"),
                       (11, "Top L Concept", CB), (12, "Top M Concept", CB), (13, "Top R Concept", CB)]:
        hdr(ws5, 2, c, txt, bg=bg, sz=8, wrap=True)

    def hash_row(ws, r, label, base, zhdr_col, zbg):
        ws.row_dimensions[r].height = 24
        sc(ws, r, 1, label, bold=True, sz=10, fc=CW, bg=zhdr_col, h="left")
        for h, cols in [('L', (2, 3, 4)), ('M', (5, 6, 7)), ('R', (8, 9, 10))]:
            hp = [p for p in base if p['hash'] == h]
            hr = [p for p in hp if p['rp'] == 'Run']; hpass = [p for p in hp if p['rp'] == 'Pass']
            n = len(hp)
            sc(ws, r, cols[0], n if n else "", sz=11, bold=True, fc="FF000000", bg=zbg, fmt="0")
            sc(ws, r, cols[1], round(len(hr) / n, 2) if n > 0 else "", sz=11, bold=True, fc="FF8B0000", bg=zbg, fmt="0%")
            sc(ws, r, cols[2], round(len(hpass) / n, 2) if n > 0 else "", sz=11, bold=True, fc="FF00008B", bg=zbg, fmt="0%")
        for col_n, h in [(11, 'L'), (12, 'M'), (13, 'R')]:
            hp = [p for p in base if p['hash'] == h]
            tc = top3(hp, 'concept', 1)
            sc(ws, r, col_n, fmt_top(tc, 0), sz=9, bg=zbg, h="left", wrap=True)

    hash_row(ws5, 3, "OVERALL", plays, CB, CL)
    for ri, zcode in enumerate(ZONE_LIST):
        hash_row(ws5, ri + 4, f"{zcode} \u2014 {ZONE_NAMES[zcode]}", [p for p in plays if p['zone'] == zcode],
                 ZONE_HDR[zcode], ZONE_BG[zcode])
    ws5.freeze_panes = "A3"

    # ── Tab 6: Down & Distance ────────────────────────────────
    ws6 = wb2.create_sheet("6. Down & Distance")
    ws6.sheet_properties.tabColor = "0E7060"; ws6.sheet_view.showGridLines = False
    NC6 = 13
    widths(ws6, [20, 8, 8, 8, 20, 20, 20, 20, 20, 20, 20, 20, 20])
    banner(ws6, 1, "DOWN & DISTANCE TENDENCIES  —  Favorite Runs, Passes & Formations by Situation", NC6, bg=CTe, sz=13, ht=28)
    for c, txt, bg in [(1, "SITUATION", CB), (2, "Plays", CB), (3, "Run%", CB), (4, "Pass%", CB),
                       (5, "#1 Run Concept", CR), (6, "#2 Run Concept", CR), (7, "#3 Run Concept", CR),
                       (8, "#1 Pass Concept", CBl), (9, "#2 Pass Concept", CBl), (10, "#3 Pass Concept", CBl),
                       (11, "#1 Formation", CPu), (12, "#2 Formation", CPu), (13, "#3 Formation", CPu)]:
        hdr(ws6, 2, c, txt, bg=bg, sz=8, wrap=True)
    for ri, (lbl, fn) in enumerate(DD_SITS):
        r = ri + 3; ws6.row_dimensions[r].height = 30
        bg = CL if ri % 2 == 0 else CW
        sp = [p for p in plays if fn(p)]
        sr = [p for p in sp if p['rp'] == 'Run']; spa = [p for p in sp if p['rp'] == 'Pass']
        sc(ws6, r, 1, lbl, bold=True, sz=10, fc=CW, bg=CTe, h="left")
        sc(ws6, r, 2, len(sp), bold=True, sz=11, fc="FF000000", bg=bg, fmt="0")
        sc(ws6, r, 3, round(len(sr) / len(sp), 2) if sp else "", bold=True, sz=12, fc="FF8B0000", bg=CRB, fmt="0%")
        sc(ws6, r, 4, round(len(spa) / len(sp), 2) if sp else "", bold=True, sz=12, fc="FF00008B", bg=CPB, fmt="0%")
        t3rc = top3_str(sr, 'concept'); t3pc = top3_str(spa, 'concept'); t3f = top3_str(sp, 'form')
        for i, cn in enumerate([5, 6, 7]): sc(ws6, r, cn, t3rc[i], sz=9, bg=CRB, wrap=True)
        for i, cn in enumerate([8, 9, 10]): sc(ws6, r, cn, t3pc[i], sz=9, bg=CPB, wrap=True)
        for i, cn in enumerate([11, 12, 13]): sc(ws6, r, cn, t3f[i], sz=9, bg="FFEDE7F6", wrap=True)
    ws6.freeze_panes = "B3"

    # ── Tab 7 / 8: Concept-level analysis (Run / Pass) ───────
    def concept_tab(ws, rp_filter, title, accent, min_n=1):
        ws.sheet_view.showGridLines = False
        NC = 13
        widths(ws, [20, 8, 8, 8, 8, 8, 8, 10, 10, 10, 10, 20, 8])
        banner(ws, 1, title, NC, bg=accent, sz=13, ht=28)
        for c, txt, bg in [(1, "CONCEPT", CB), (2, "Called", CB), (3, "Avg Yd", CB), (4, "Expl%", CB),
                           (5, "Succ%", CB), (6, "Dir R%", accent), (7, "Dir L%", accent),
                           (8, "1st Dn\nSucc", accent), (9, "2nd Dn\nSucc", accent), (10, "3rd Dn\nSucc", accent),
                           (11, "RedZone\nSucc", accent), (12, "Top Form", CB), (13, "Top Hash", CB)]:
            hdr(ws, 2, c, txt, bg=bg, sz=8, wrap=True)
        grp = {}
        for p in plays:
            if p['rp'] != rp_filter: continue
            c = str(p['concept']).strip()
            if c in ('', 'nan', 'None'): continue
            grp.setdefault(c, []).append(p)
        ranked = sorted(grp.items(), key=lambda kv: -len(kv[1]))
        ranked = [(c, g) for c, g in ranked if len(g) >= min_n]
        for ri, (concept, g) in enumerate(ranked):
            r = ri + 3; ws.row_dimensions[r].height = 22
            bg = CL if ri % 2 == 0 else CW
            small = " *" if len(g) < 5 else ""
            expl_pct = len([p for p in g if p['expl']]) / len(g)
            sr = _sr(g)
            rd = [p for p in g if p['dir'] == 'R']; ld = [p for p in g if p['dir'] == 'L']
            dknown = len(rd) + len(ld)
            sc(ws, r, 1, concept + small, bold=True, sz=9, fc=CW, bg=accent, h="left")
            sc(ws, r, 2, len(g), bold=True, sz=10, fc="FF000000", bg=bg, fmt="0")
            sc(ws, r, 3, round(_avg(g), 1), bold=True, sz=10, fc="FF0E7060", bg="FFE8F8E8", fmt="0.0")
            sc(ws, r, 4, round(expl_pct, 2), sz=9, fc="FF0E7060", bg="FFE8F8E8", fmt="0%")
            sc(ws, r, 5, (round(sr / 100, 2) if sr is not None else "—"), bold=True, sz=10, fc="FF0E7060",
               bg="FFE8F8E8", fmt="0%" if sr is not None else "General")
            sc(ws, r, 6, round(len(rd) / dknown, 2) if dknown else "—", sz=9, fc="FF8B0000", bg=CRB,
               fmt="0%" if dknown else "General")
            sc(ws, r, 7, round(len(ld) / dknown, 2) if dknown else "—", sz=9, fc="FF00008B", bg=CPB,
               fmt="0%" if dknown else "General")
            for dn, cn in [(1, 8), (2, 9), (3, 10)]:
                dg = [p for p in g if p['dn'] == dn]
                dsr = _sr(dg)
                sc(ws, r, cn, (round(dsr / 100, 2) if dsr is not None else "—"), sz=9, fc="FF4A235A",
                   bg="FFEDE7F6", fmt="0%" if dsr is not None else "General")
            rzg = [p for p in g if p['zone'] in ('RZ', 'GL')]
            rzsr = _sr(rzg)
            sc(ws, r, 11, (round(rzsr / 100, 2) if rzsr is not None else "—"), sz=9, fc="FFC0392B",
               bg=CRB, fmt="0%" if rzsr is not None else "General")
            tf = top3(g, 'form', 1); th = top3(g, 'hash', 1)
            sc(ws, r, 12, tf[0]['v'] if tf else "—", sz=8, bg=bg, h="left", wrap=True)
            sc(ws, r, 13, th[0]['v'] if th else "—", sz=9, bg=bg)
        if not ranked:
            ws.merge_cells(f"A3:{gcl(NC)}3")
            c = ws.cell(row=3, column=1, value="Not enough tagged data — check the CONCEPT column is filled in.")
            c.font = Font(name=FN, sz=10, italic=True, color=CDG); c.alignment = Alignment(horizontal="center")
        fr = len(ranked) + 4
        ws.cell(row=fr, column=1, value="* = small sample (under 5 calls) — read with caution").font = Font(
            name=FN, sz=8, italic=True, color=CDG)
        ws.freeze_panes = "B3"

    ws7 = wb2.create_sheet("7. Run Concepts"); ws7.sheet_properties.tabColor = "C0392B"
    concept_tab(ws7, "Run", "RUN CONCEPTS  \u2014  Explosive (10+) & Success Rate by Concept (All Calls)", "FFC0392B")

    ws8 = wb2.create_sheet("8. Pass Concepts"); ws8.sheet_properties.tabColor = "1A5276"
    concept_tab(ws8, "Pass", "PASS CONCEPTS  \u2014  Explosive (15+) & Success Rate by Concept (All Calls)", "FF1A5276")

    # ── Tab 9: Formation Tendencies ───────────────────────────
    ws9 = wb2.create_sheet("9. Formation Tendencies")
    ws9.sheet_properties.tabColor = "4A235A"; ws9.sheet_view.showGridLines = False
    NC9 = 12
    widths(ws9, [22, 8, 8, 8, 8, 8, 18, 18, 18, 18, 16, 14])
    banner(ws9, 1, "FORMATION TENDENCIES  \u2014  Favorite Runs & Passes by Formation (3+ snaps)", NC9, bg=CPu, sz=13, ht=28)
    for c, txt, bg in [(1, "FORMATION", CB), (2, "Snaps", CB), (3, "Run%", CB), (4, "Pass%", CB),
                       (5, "Run R%", CR), (6, "Run L%", CR), (7, "#1 Run Concept", CR), (8, "#2 Run Concept", CR),
                       (9, "#1 Pass Concept", CBl), (10, "#2 Pass Concept", CBl),
                       (11, "Top Down/Dist", CPu), (12, "Top Zone", CPu)]:
        hdr(ws9, 2, c, txt, bg=bg, sz=8, wrap=True)
    fgroups = {}
    for p in plays:
        f = p['form']
        if f.strip() in ('', 'nan', 'None'): continue
        fgroups.setdefault(f, []).append(p)
    ranked = sorted(fgroups.items(), key=lambda kv: -len(kv[1]))
    ranked = [(f, g) for f, g in ranked if len(g) >= 3]
    for ri, (f, g) in enumerate(ranked):
        r = ri + 3; ws9.row_dimensions[r].height = 24
        bg = CL if ri % 2 == 0 else CW
        gr = [p for p in g if p['rp'] == 'Run']; gp = [p for p in g if p['rp'] == 'Pass']
        rd = [p for p in gr if p['dir'] == 'R']; ld = [p for p in gr if p['dir'] == 'L']
        dknown = len(rd) + len(ld)
        sc(ws9, r, 1, f, bold=True, sz=9, fc=CW, bg=CPu, h="left")
        sc(ws9, r, 2, len(g), bold=True, sz=10, fc="FF000000", bg=bg, fmt="0")
        sc(ws9, r, 3, round(len(gr) / len(g), 2), bold=True, sz=10, fc="FF8B0000", bg=CRB, fmt="0%")
        sc(ws9, r, 4, round(len(gp) / len(g), 2), bold=True, sz=10, fc="FF00008B", bg=CPB, fmt="0%")
        sc(ws9, r, 5, round(len(rd) / dknown, 2) if dknown else "—", sz=9, bg=bg, fmt="0%" if dknown else "General")
        sc(ws9, r, 6, round(len(ld) / dknown, 2) if dknown else "—", sz=9, bg=bg, fmt="0%" if dknown else "General")
        t3rc = top3_str(gr, 'concept', 2); t3pc = top3_str(gp, 'concept', 2)
        for i, cn in enumerate([7, 8]): sc(ws9, r, cn, t3rc[i], sz=8, bg=CRB, wrap=True)
        for i, cn in enumerate([9, 10]): sc(ws9, r, cn, t3pc[i], sz=8, bg=CPB, wrap=True)
        dd_top = Counter(dd_bucket(p) for p in g).most_common(1)
        sc(ws9, r, 11, f"{dd_top[0][0]} ({dd_top[0][1]})" if dd_top else "—", sz=8, bg=bg, h="left", wrap=True)
        z_top = Counter(ZONE_NAMES[p['zone']].split("  ")[0] for p in g).most_common(1)
        sc(ws9, r, 12, f"{z_top[0][0]} ({z_top[0][1]})" if z_top else "—", sz=8, bg=bg, h="left", wrap=True)
    if not ranked:
        ws9.merge_cells(f"A3:{gcl(NC9)}3")
        c = ws9.cell(row=3, column=1, value="Not enough formation data tagged.")
        c.font = Font(name=FN, sz=10, italic=True, color=CDG); c.alignment = Alignment(horizontal="center")
    ws9.freeze_panes = "B3"

    # ── Tab 10: Situational Summary ───────────────────────────
    ws10 = wb2.create_sheet("10. Situational Summary")
    ws10.sheet_properties.tabColor = "F1C40F"; ws10.sheet_view.showGridLines = False
    NC10 = 10
    widths(ws10, [18, 9, 9, 10, 10, 10, 20, 20, 22, 26])
    banner(ws10, 1, "SITUATIONAL SUMMARY", NC10, bg="FFF1C40F", fc=CB, sz=13, ht=28)
    s10h = ["Situation", "Run\nCount", "Pass\nCount", "L Hash\nRun%", "M Hash\nRun%", "R Hash\nRun%",
            "Top Run Concept", "Top Pass Concept", "Best Call", "Notes"]
    for ci, h in enumerate(s10h): hdr(ws10, 2, ci + 1, h, bg=CB, sz=9, wrap=True)

    def sit(dn=None, dmin=None, dmax=None, zone=None, custom=None):
        out = []
        for p in plays:
            if custom is not None:
                if custom(p): out.append(p)
                continue
            if dn and p['dn'] != dn: continue
            if dmin and p['dist'] < dmin: continue
            if dmax and p['dist'] > dmax: continue
            if zone and p['zone'] != zone: continue
            out.append(p)
        return out

    sits10 = [
        ("1ST DOWN", dict(dn=1)),
        ("2ND & LONG", dict(dn=2, dmin=7)),
        ("2ND & MEDIUM", dict(dn=2, dmin=4, dmax=6)),
        ("2ND & SHORT", dict(dn=2, dmax=3)),
        ("3RD & LONG", dict(dn=3, dmin=7)),
        ("3RD & MEDIUM", dict(dn=3, dmin=4, dmax=6)),
        ("3RD & SHORT", dict(dn=3, dmax=3)),
        ("4TH DOWN", dict(dn=4)),
        ("RED ZONE", dict(zone="RZ")),
        ("GOAL LINE", dict(zone="GL")),
        ("BACKED UP", dict(zone="BZ")),
        ("COMING OUT", dict(zone="OF")),
        ("TWO-MINUTE", dict(custom=lambda p: p['dn'] == 3 and p['dist'] >= 5)),
        ("MUST-HAVE", dict(dn=4)),
    ]
    sit_colors = ["FF0E7060", "FF1A5276", "FF1A5276", "FF1A5276",
                  "FFC0392B", "FFC0392B", "FFC0392B", "FF7B241C",
                  "FFC0392B", "FF4A235A", "FF0E7060", "FF0E7060", "FF7D6608", "FF16213E"]
    for ri, ((lbl, args), color) in enumerate(zip(sits10, sit_colors)):
        r = ri + 3; ws10.row_dimensions[r].height = 30
        sc(ws10, r, 1, lbl, bold=True, sz=9, fc=CW, bg=color, h="left")
        sp = sit(**args)
        sr = [p for p in sp if p['rp'] == 'Run']; spass = [p for p in sp if p['rp'] == 'Pass']
        l_p = [p for p in sp if p['hash'] == 'L']; m_p = [p for p in sp if p['hash'] == 'M']; r_p = [p for p in sp if p['hash'] == 'R']
        l_r = len([p for p in l_p if p['rp'] == 'Run']); m_r = len([p for p in m_p if p['rp'] == 'Run']); r_r = len([p for p in r_p if p['rp'] == 'Run'])
        tf = top3(sr, 'concept', 1); tc = top3(spass, 'concept', 1)
        sc(ws10, r, 2, len(sr), bold=True, sz=12, fc="FF8B0000", bg="FFFDE8E8", fmt="0")
        sc(ws10, r, 3, len(spass), bold=True, sz=12, fc="FF00008B", bg="FFE8F0FE", fmt="0")
        sc(ws10, r, 4, round(l_r / len(l_p), 2) if l_p else "", sz=10, fc="FF6C3483", bg="FFEAF0FF", fmt="0%")
        sc(ws10, r, 5, round(m_r / len(m_p), 2) if m_p else "", sz=10, fc="FF1A5276", bg="FFE8F0FE", fmt="0%")
        sc(ws10, r, 6, round(r_r / len(r_p), 2) if r_p else "", sz=10, fc="FF784212", bg="FFFFF0EA", fmt="0%")
        sc(ws10, r, 7, fmt_top(tf, 0), sz=9, bg="FFFDE8E8", wrap=True, h="left")
        sc(ws10, r, 8, fmt_top(tc, 0), sz=9, bg="FFE8F0FE", wrap=True, h="left")
        sc(ws10, r, 9, "", bg=CYB, sz=9, wrap=True, v="top")
        sc(ws10, r, 10, "", bg=CL if ri % 2 == 0 else CW, sz=9, wrap=True, v="top")
    ws10.freeze_panes = "D3"

    # ── Tab 11: Call Sheet Builder ─────────────────────────────
    ws11 = wb2.create_sheet("11. Call Sheet Builder")
    ws11.sheet_properties.tabColor = "0E7060"; ws11.sheet_view.showGridLines = False
    NC11 = 11
    widths(ws11, [20, 10, 12, 26, 18, 14, 14, 14, 14, 9, 24])
    banner(ws11, 1, "CALL SHEET BUILDER  \u2014  Opponent tendency auto-filled \u00b7 Yellow = your calls", NC11, bg=CTe, sz=13, ht=28)
    c11h = ["Situation", "Field\nZone", "Down /\nDistance", "Opponent Tendency", "Formation / Alert",
            "Best Pressure", "Best Coverage", "Best Front", "Adjustment", "Priority", "Notes"]
    for ci, h in enumerate(c11h): hdr(ws11, 2, ci + 1, h, bg=CTe, sz=9, wrap=True)
    cs_filters = [
        ("1st & 10", lambda p: p['dn'] == 1 and p['dist'] >= 8, "", ""),
        ("1st & 10 (Own Half)", lambda p: p['dn'] == 1 and p['dist'] >= 8 and p['zone'] in ('BZ', 'OF'), "Own Half", ""),
        ("1st & 10 (Opp Half)", lambda p: p['dn'] == 1 and p['dist'] >= 8 and p['zone'] in ('MF', 'FZ'), "Opp Half", ""),
        ("2nd & Long (8+)", lambda p: p['dn'] == 2 and p['dist'] >= 8, "", ""),
        ("2nd & Medium (4-7)", lambda p: p['dn'] == 2 and 4 <= p['dist'] <= 7, "", ""),
        ("2nd & Short (1-3)", lambda p: p['dn'] == 2 and p['dist'] <= 3, "", ""),
        ("3rd & Long (7+)", lambda p: p['dn'] == 3 and p['dist'] >= 7, "", ""),
        ("3rd & Medium (4-6)", lambda p: p['dn'] == 3 and 4 <= p['dist'] <= 6, "", ""),
        ("3rd & Short (1-3)", lambda p: p['dn'] == 3 and p['dist'] <= 3, "", ""),
        ("4th Down", lambda p: p['dn'] == 4, "", ""),
        ("Red Zone \u2014 1st", lambda p: p['zone'] == 'RZ' and p['dn'] == 1, "Red Zone", ""),
        ("Red Zone \u2014 2nd", lambda p: p['zone'] == 'RZ' and p['dn'] == 2, "Red Zone", ""),
        ("Red Zone \u2014 3rd", lambda p: p['zone'] == 'RZ' and p['dn'] == 3, "Red Zone", ""),
        ("Goal Line", lambda p: p['zone'] == 'GL', "Goal Line", ""),
        ("Backed Up", lambda p: p['zone'] == 'BZ', "Backed Up", ""),
        ("Coming Out", lambda p: p['zone'] == 'OF', "Open Field", ""),
        ("Two-Minute (Lead)", None, "", ""),
        ("Two-Minute (Trail)", None, "", ""),
    ]
    r = 3
    for ri, (lbl, fn, zonelbl, _n) in enumerate(cs_filters):
        ws11.row_dimensions[r].height = 24
        bg = "FFF0FFF0" if ri % 2 == 0 else CW
        sp = [p for p in plays if fn(p)] if fn else []
        sr = [p for p in sp if p['rp'] == 'Run']; spa = [p for p in sp if p['rp'] == 'Pass']
        tend = ""
        if sp:
            tend = f"{pct(len(sr), len(sp))}% Run / {pct(len(spa), len(sp))}% Pass  ({len(sp)} snaps)"
        alertf = top3(sp, 'form', 1)
        alert_txt = fmt_top(alertf, 0) if sp else ""
        sc(ws11, r, 1, lbl, bold=True, sz=9, fc=CW, bg=CGr, h="left")
        sc(ws11, r, 2, zonelbl, sz=9, bg=bg)
        sc(ws11, r, 3, "", sz=9, bg=bg)
        sc(ws11, r, 4, tend, sz=9, bg=bg, h="left", wrap=True)
        sc(ws11, r, 5, alert_txt, sz=9, bg=bg, h="left", wrap=True)
        for ci in range(6, NC11 + 1):
            sc(ws11, r, ci, "", bg=CYB, sz=9, wrap=True, v="top")
        r += 1
    ws11.freeze_panes = "B3"

    # ── Tab 12: Game Day Call Sheet (printable one-pager) ─────
    ws12 = wb2.create_sheet("12. Game Day Call Sheet")
    ws12.sheet_properties.tabColor = "16213E"; ws12.sheet_view.showGridLines = False
    ws12.page_setup.orientation = "landscape"
    ws12.page_setup.fitToPage = True; ws12.page_setup.fitToWidth = 1; ws12.page_setup.fitToHeight = 1
    widths(ws12, [16, 10, 10, 18, 20, 20, 3, 10, 30, 3, 3])
    heading = f"GAME DAY CALL SHEET   \u00b7   {opp or 'Opponent'}" + (f"   \u00b7   WK {week}" if week else "")
    banner(ws12, 1, heading, 11, bg=CB, sz=14, ht=30)

    # left block: down & distance
    r = 3
    ws12.merge_cells(start_row=r, start_column=1, end_row=r, end_column=6)
    c = ws12.cell(row=r, column=1, value="DOWN & DISTANCE")
    c.font = Font(name=FN, bold=True, sz=11, color=CW); c.fill = fil(CTe)
    c.alignment = Alignment(horizontal="center"); ws12.row_dimensions[r].height = 18
    r += 1
    for ci, h in enumerate(["Situation", "Run%", "Pass%", "Top Run", "Top Pass", "Form"], 1):
        hdr(ws12, r, ci, h, bg=CB, sz=8)
    r += 1
    dd_top_sits = [("1st & 10", DD_SITS[0][1]), ("2nd & Short", DD_SITS[4][1]), ("2nd & Med", DD_SITS[3][1]),
                   ("2nd & Long", DD_SITS[2][1]), ("3rd & Short", DD_SITS[7][1]), ("3rd & Med", DD_SITS[6][1]),
                   ("3rd & Long", DD_SITS[5][1]), ("4th Down", DD_SITS[8][1])]
    for ri, (lbl, fn) in enumerate(dd_top_sits):
        bg = CL if ri % 2 == 0 else CW
        sp = [p for p in plays if fn(p)]
        sr = [p for p in sp if p['rp'] == 'Run']; spa = [p for p in sp if p['rp'] == 'Pass']
        sc(ws12, r, 1, lbl, bold=True, sz=9, fc=CB, bg=bg, h="left")
        sc(ws12, r, 2, round(len(sr) / len(sp), 2) if sp else "", sz=9, fc="FF8B0000", bg=bg, fmt="0%")
        sc(ws12, r, 3, round(len(spa) / len(sp), 2) if sp else "", sz=9, fc="FF00008B", bg=bg, fmt="0%")
        sc(ws12, r, 4, fmt_top(top3(sr, 'concept', 1), 0), sz=8, bg=bg, h="left")
        sc(ws12, r, 5, fmt_top(top3(spa, 'concept', 1), 0), sz=8, bg=bg, h="left")
        sc(ws12, r, 6, fmt_top(top3(sp, 'form', 1), 0), sz=8, bg=bg, h="left")
        ws12.row_dimensions[r].height = 15
        r += 1

    r += 1
    ws12.merge_cells(start_row=r, start_column=1, end_row=r, end_column=6)
    c = ws12.cell(row=r, column=1, value="FIELD ZONE")
    c.font = Font(name=FN, bold=True, sz=11, color=CW); c.fill = fil(CTe)
    c.alignment = Alignment(horizontal="center"); ws12.row_dimensions[r].height = 18
    r += 1
    for ci, h in enumerate(["Zone", "Run%", "Pass%", "Top Run", "Top Pass", "Form"], 1):
        hdr(ws12, r, ci, h, bg=CB, sz=8)
    r += 1
    zone_disp = [("Coming Out", "OF"), ("Midfield", "MF"), ("Fringe", "FZ"), ("Red Zone", "RZ"), ("Goal Line", "GL"), ("Backed Up", "BZ")]
    for ri, (lbl, zc) in enumerate(zone_disp):
        bg = CL if ri % 2 == 0 else CW
        sp = [p for p in plays if p['zone'] == zc]
        sr = [p for p in sp if p['rp'] == 'Run']; spa = [p for p in sp if p['rp'] == 'Pass']
        sc(ws12, r, 1, lbl, bold=True, sz=9, fc=CB, bg=bg, h="left")
        sc(ws12, r, 2, round(len(sr) / len(sp), 2) if sp else "", sz=9, fc="FF8B0000", bg=bg, fmt="0%")
        sc(ws12, r, 3, round(len(spa) / len(sp), 2) if sp else "", sz=9, fc="FF00008B", bg=bg, fmt="0%")
        sc(ws12, r, 4, fmt_top(top3(sr, 'concept', 1), 0), sz=8, bg=bg, h="left")
        sc(ws12, r, 5, fmt_top(top3(spa, 'concept', 1), 0), sz=8, bg=bg, h="left")
        sc(ws12, r, 6, fmt_top(top3(sp, 'form', 1), 0), sz=8, bg=bg, h="left")
        ws12.row_dimensions[r].height = 15
        r += 1

    # right block: biggest tendencies / formation alerts / heavy pass
    r2 = 3
    ws12.merge_cells(start_row=r2, start_column=8, end_row=r2, end_column=9)
    c = ws12.cell(row=r2, column=8, value="BIGGEST TENDENCIES")
    c.font = Font(name=FN, bold=True, sz=11, color=CW); c.fill = fil(CR)
    c.alignment = Alignment(horizontal="center"); ws12.row_dimensions[r2].height = 18
    r2 += 1
    big = compute_biggest_tendencies(plays, top_n=6)
    for i, txt in enumerate(big):
        bg = CL if i % 2 == 0 else CW
        sc(ws12, r2, 8, f"#{i+1}", bold=True, sz=9, fc=CW, bg=CR)
        sc(ws12, r2, 9, txt, sz=9, bg=bg, h="left", wrap=True)
        ws12.row_dimensions[r2].height = 16
        r2 += 1
    if not big:
        sc(ws12, r2, 8, "", bg=CW)
        sc(ws12, r2, 9, "Not enough tagged data for alerts yet.", sz=9)
        r2 += 1
    r2 += 1
    ws12.merge_cells(start_row=r2, start_column=8, end_row=r2, end_column=9)
    c = ws12.cell(row=r2, column=8, value="FORMATION ALERTS")
    c.font = Font(name=FN, bold=True, sz=11, color=CW); c.fill = fil(CPu)
    c.alignment = Alignment(horizontal="center"); ws12.row_dimensions[r2].height = 18
    r2 += 1
    for f, tag in compute_formation_alerts(plays, top_n=6):
        bg = CL if (r2 % 2 == 0) else CW
        sc(ws12, r2, 8, f, bold=True, sz=9, fc=CW, bg=CPu, h="left")
        sc(ws12, r2, 9, tag, sz=9, bg=bg, h="left", wrap=True)
        ws12.row_dimensions[r2].height = 16
        r2 += 1
    r2 += 1
    ws12.merge_cells(start_row=r2, start_column=8, end_row=r2, end_column=9)
    c = ws12.cell(row=r2, column=8, value="HEAVY PASS SITUATIONS")
    c.font = Font(name=FN, bold=True, sz=11, color=CW); c.fill = fil(CBl)
    c.alignment = Alignment(horizontal="center"); ws12.row_dimensions[r2].height = 18
    r2 += 1
    heavy = compute_heavy_pass_situations(plays)
    for i, txt in enumerate(heavy):
        bg = CL if i % 2 == 0 else CW
        sc(ws12, r2, 8, f"#{i+1}", bold=True, sz=9, fc=CW, bg=CBl)
        sc(ws12, r2, 9, txt, sz=9, bg=bg, h="left", wrap=True)
        ws12.row_dimensions[r2].height = 16
        r2 += 1
    ws12.freeze_panes = "A4"

    # ── Tab 13 / 14: Group Tendencies (FORM FAMILY / FIB) ─────
    def group_tab(ws, group_key, title, accent, tab_color, empty_label="(Blank)", min_n=1):
        ws.sheet_properties.tabColor = tab_color
        ws.sheet_view.showGridLines = False
        NC = 13
        widths(ws, [20, 8, 8, 8, 18, 18, 18, 18, 18, 18, 16, 16, 16])
        banner(ws, 1, title, NC, bg=accent, sz=13, ht=28)
        for c, txt, bg in [(1, "GROUP", CB), (2, "Snaps", CB), (3, "Run%", CB), (4, "Pass%", CB),
                           (5, "#1 Run Play", CR), (6, "#2 Run Play", CR), (7, "#3 Run Play", CR),
                           (8, "#1 Pass Play", CBl), (9, "#2 Pass Play", CBl), (10, "#3 Pass Play", CBl),
                           (11, "#1 Formation", CPu), (12, "#2 Formation", CPu), (13, "#3 Formation", CPu)]:
            hdr(ws, 2, c, txt, bg=bg, sz=8, wrap=True)
        groups = {}
        for p in plays:
            v = str(p.get(group_key, '')).strip()
            if v in ('', 'nan', 'None'): v = empty_label
            groups.setdefault(v, []).append(p)
        ranked = sorted(groups.items(), key=lambda kv: -len(kv[1]))
        ranked = [(v, g) for v, g in ranked if len(g) >= min_n]
        for ri, (v, g) in enumerate(ranked):
            r = ri + 3; ws.row_dimensions[r].height = 24
            bg = CL if ri % 2 == 0 else CW
            gr = [p for p in g if p['rp'] == 'Run']; gp = [p for p in g if p['rp'] == 'Pass']
            sc(ws, r, 1, v, bold=True, sz=9, fc=CW, bg=accent, h="left")
            sc(ws, r, 2, len(g), bold=True, sz=10, fc="FF000000", bg=bg, fmt="0")
            sc(ws, r, 3, round(len(gr) / len(g), 2) if g else "", bold=True, sz=10, fc="FF8B0000", bg=CRB, fmt="0%")
            sc(ws, r, 4, round(len(gp) / len(g), 2) if g else "", bold=True, sz=10, fc="FF00008B", bg=CPB, fmt="0%")
            t3rc = top3_str(gr, 'concept', 3); t3pc = top3_str(gp, 'concept', 3); t3f = top3_str(g, 'form', 3)
            for i, cn in enumerate([5, 6, 7]): sc(ws, r, cn, t3rc[i], sz=9, bg=CRB, wrap=True)
            for i, cn in enumerate([8, 9, 10]): sc(ws, r, cn, t3pc[i], sz=9, bg=CPB, wrap=True)
            for i, cn in enumerate([11, 12, 13]): sc(ws, r, cn, t3f[i], sz=9, bg="FFEDE7F6", wrap=True)
        if not ranked:
            ws.merge_cells(f"A3:{gcl(NC)}3")
            c = ws.cell(row=3, column=1, value="Not enough tagged data for this section.")
            c.font = Font(name=FN, sz=10, italic=True, color=CDG); c.alignment = Alignment(horizontal="center")
        ws.freeze_panes = "B3"

    ws13 = wb2.create_sheet("13. Form Family Tendencies")
    group_tab(ws13, 'form_family', "FORM FAMILY TENDENCIES  \u2014  Run/Pass Split, Favorite Plays & Formations",
              "FF6C3483", "6C3483")

    ws14 = wb2.create_sheet("14. FIB Tendencies")
    group_tab(ws14, 'fib', "FIB TENDENCIES  \u2014  Run/Pass Split, Favorite Plays & Formations",
              "FF784212", "784212", empty_label="Not FIB")

    ws14b = wb2.create_sheet("15. Back Depth")
    group_tab(ws14b, 'back_depth', "BACK DEPTH TENDENCIES  \u2014  Run/Pass Split, Favorite Plays & Formations",
              "FF0E7060", "0E7060", empty_label="(Blank)")

    # ── Tab 16: Open/Closed (with cross-break by Form Family) ─
    ws14c = wb2.create_sheet("16. Open-Closed")

    def build_open_closed_tab(ws):
        ws.sheet_properties.tabColor = "4A235A"
        ws.sheet_view.showGridLines = False
        NC = 13
        widths(ws, [20, 8, 8, 8, 18, 18, 18, 18, 18, 18, 16, 16, 16])
        banner(ws, 1, "OPEN / CLOSED TENDENCIES  \u2014  Run/Pass Split, Favorite Plays & Formations",
               NC, bg="FF4A235A", sz=13, ht=28)
        for c, txt, bg in [(1, "GROUP", CB), (2, "Snaps", CB), (3, "Run%", CB), (4, "Pass%", CB),
                           (5, "#1 Run Play", CR), (6, "#2 Run Play", CR), (7, "#3 Run Play", CR),
                           (8, "#1 Pass Play", CBl), (9, "#2 Pass Play", CBl), (10, "#3 Pass Play", CBl),
                           (11, "#1 Formation", CPu), (12, "#2 Formation", CPu), (13, "#3 Formation", CPu)]:
            hdr(ws, 2, c, txt, bg=bg, sz=8, wrap=True)
        groups = {}
        for p in plays:
            v = str(p.get('open_close', '')).strip()
            if v in ('', 'nan', 'None'): v = "(Blank)"
            groups.setdefault(v, []).append(p)
        ranked = sorted(groups.items(), key=lambda kv: -len(kv[1]))
        row = 3
        for ri, (v, g) in enumerate(ranked):
            bg = CL if ri % 2 == 0 else CW
            gr = [p for p in g if p['rp'] == 'Run']; gp = [p for p in g if p['rp'] == 'Pass']
            sc(ws, row, 1, v, bold=True, sz=9, fc=CW, bg="FF4A235A", h="left")
            sc(ws, row, 2, len(g), bold=True, sz=10, fc="FF000000", bg=bg, fmt="0")
            sc(ws, row, 3, round(len(gr) / len(g), 2) if g else "", bold=True, sz=10, fc="FF8B0000", bg=CRB, fmt="0%")
            sc(ws, row, 4, round(len(gp) / len(g), 2) if g else "", bold=True, sz=10, fc="FF00008B", bg=CPB, fmt="0%")
            t3rc = top3_str(gr, 'concept', 3); t3pc = top3_str(gp, 'concept', 3); t3f = top3_str(g, 'form', 3)
            for i, cn in enumerate([5, 6, 7]): sc(ws, row, cn, t3rc[i], sz=9, bg=CRB, wrap=True)
            for i, cn in enumerate([8, 9, 10]): sc(ws, row, cn, t3pc[i], sz=9, bg=CPB, wrap=True)
            for i, cn in enumerate([11, 12, 13]): sc(ws, row, cn, t3f[i], sz=9, bg="FFEDE7F6", wrap=True)
            row += 1
        if not ranked:
            ws.merge_cells(f"A3:{gcl(NC)}3")
            c = ws.cell(row=3, column=1, value="Not enough tagged data for this section.")
            c.font = Font(name=FN, sz=10, italic=True, color=CDG); c.alignment = Alignment(horizontal="center")
            row = 4
        ws.freeze_panes = "B3"

        # ── Cross-break: Open/Closed within each Formation Family ──
        row += 1
        banner(ws, row, "OPEN / CLOSED BY FORMATION FAMILY", NC, bg="FF6C3483", sz=12, ht=26)
        row += 1
        for c, txt, bg in [(1, "FAMILY / GROUP", CB), (2, "Snaps", CB), (3, "Run%", CB), (4, "Pass%", CB),
                           (5, "#1 Run Play", CR), (6, "#2 Run Play", CR), (7, "#3 Run Play", CR),
                           (8, "#1 Pass Play", CBl), (9, "#2 Pass Play", CBl), (10, "#3 Pass Play", CBl),
                           (11, "#1 Formation", CPu), (12, "#2 Formation", CPu), (13, "#3 Formation", CPu)]:
            hdr(ws, row, c, txt, bg=bg, sz=8, wrap=True)
        row += 1
        fam_groups = {}
        for p in plays:
            fam = str(p.get('form_family', '')).strip()
            if fam in ('', 'nan', 'None'): continue
            fam_groups.setdefault(fam, []).append(p)
        fam_ranked = sorted(fam_groups.items(), key=lambda kv: -len(kv[1]))
        combo_i = 0
        for fam, fam_plays in fam_ranked:
            ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=NC)
            c = ws.cell(row=row, column=1, value=f"  {fam}  ({len(fam_plays)} snaps)")
            c.font = Font(name=FN, bold=True, size=10, color=CW)
            c.fill = fil("FF6C3483")
            c.alignment = Alignment(horizontal="left", vertical="center")
            ws.row_dimensions[row].height = 18
            row += 1
            oc_groups = {}
            for p in fam_plays:
                v = str(p.get('open_close', '')).strip()
                if v in ('', 'nan', 'None'): v = "(Blank)"
                oc_groups.setdefault(v, []).append(p)
            oc_ranked = sorted(oc_groups.items(), key=lambda kv: -len(kv[1]))
            for v, g in oc_ranked:
                bg = CL if combo_i % 2 == 0 else CW
                combo_i += 1
                gr = [p for p in g if p['rp'] == 'Run']; gp = [p for p in g if p['rp'] == 'Pass']
                sc(ws, row, 1, f"    {v}", bold=True, sz=9, fc="FF4A235A", bg=bg, h="left")
                sc(ws, row, 2, len(g), bold=True, sz=9, fc="FF000000", bg=bg, fmt="0")
                sc(ws, row, 3, round(len(gr) / len(g), 2) if g else "", sz=9, fc="FF8B0000", bg=CRB, fmt="0%")
                sc(ws, row, 4, round(len(gp) / len(g), 2) if g else "", sz=9, fc="FF00008B", bg=CPB, fmt="0%")
                t3rc = top3_str(gr, 'concept', 3); t3pc = top3_str(gp, 'concept', 3); t3f = top3_str(g, 'form', 3)
                for i, cn in enumerate([5, 6, 7]): sc(ws, row, cn, t3rc[i], sz=8, bg=CRB, wrap=True)
                for i, cn in enumerate([8, 9, 10]): sc(ws, row, cn, t3pc[i], sz=8, bg=CPB, wrap=True)
                for i, cn in enumerate([11, 12, 13]): sc(ws, row, cn, t3f[i], sz=8, bg="FFEDE7F6", wrap=True)
                row += 1
        if not fam_ranked:
            ws.merge_cells(f"A{row}:{gcl(NC)}{row}")
            c = ws.cell(row=row, column=1, value="Not enough Form Family data tagged.")
            c.font = Font(name=FN, sz=10, italic=True, color=CDG); c.alignment = Alignment(horizontal="center")

    build_open_closed_tab(ws14c)

    # ── Tab 17: Stats (Passing / Rushing / Receiving) ──────────
    ws15 = wb2.create_sheet("17. Stats")
    ws15.sheet_properties.tabColor = "16213E"; ws15.sheet_view.showGridLines = False
    NC15 = 11
    widths(ws15, [10, 8, 8, 8, 8, 8, 8, 8, 8, 8, 8])

    passing, rushing, receiving = compute_player_stats(plays)

    def _dash(n):
        return n if n else "—"

    row = 1
    banner(ws15, row, "PASSING", NC15, bg=CR, sz=13, ht=26); row += 1
    for c, txt in enumerate(["Player", "Cmp", "Att", "%", "Yds", "Y/C", "Lng", "TD", "Fum", "Int", "Rat"], 1):
        hdr(ws15, row, c, txt, bg=CB, sz=9)
    row += 1
    pass_ranked = sorted(passing.items(), key=lambda kv: -kv[1]['yds'])
    for ri, (player, d) in enumerate(pass_ranked):
        bg = CL if ri % 2 == 0 else CW
        att, cmp_, yds, td, intc, fum = d['att'], d['cmp'], d['yds'], d['td'], d['int'], d['fum']
        rating = (8.4 * yds + 330 * td + 100 * cmp_ - 200 * intc) / att if att else 0
        sc(ws15, row, 1, player, bold=True, sz=10, fc=CW, bg=CBl, h="left")
        sc(ws15, row, 2, cmp_, sz=9, bg=bg, fmt="0")
        sc(ws15, row, 3, att, sz=9, bg=bg, fmt="0")
        sc(ws15, row, 4, round(cmp_ / att, 3) if att else "", sz=9, bg=bg, fmt="0.0%")
        sc(ws15, row, 5, yds, bold=True, sz=9, fc="FF0E7060", bg="FFE8F8E8", fmt="0")
        sc(ws15, row, 6, round(yds / cmp_, 1) if cmp_ else "", sz=9, bg=bg, fmt="0.0")
        sc(ws15, row, 7, d['lng'] if d['lng'] is not None else "", sz=9, bg=bg, fmt="0")
        sc(ws15, row, 8, _dash(td), sz=9, bg=bg, fmt="0" if td else "General")
        sc(ws15, row, 9, _dash(fum), sz=9, bg=bg, fmt="0" if fum else "General")
        sc(ws15, row, 10, _dash(intc), sz=9, bg=bg, fmt="0" if intc else "General")
        sc(ws15, row, 11, round(rating, 1), bold=True, sz=9, fc="FFC0392B", bg=CRB, fmt="0.0")
        row += 1
    if not pass_ranked:
        ws15.merge_cells(start_row=row, start_column=1, end_row=row, end_column=NC15)
        c = ws15.cell(row=row, column=1, value="No passer data tagged (OPP PASSER column).")
        c.font = Font(name=FN, sz=9, italic=True, color=CDG); c.alignment = Alignment(horizontal="center")
        row += 1
    row += 1

    banner(ws15, row, "RUSHING", NC15, bg=CR, sz=13, ht=26); row += 1
    for c, txt in enumerate(["Player", "Att", "Yds", "Avg", "Lng", "TD", "Fum"], 1):
        hdr(ws15, row, c, txt, bg=CB, sz=9)
    row += 1
    rush_ranked = sorted(rushing.items(), key=lambda kv: -kv[1]['yds'])
    for ri, (player, d) in enumerate(rush_ranked):
        bg = CL if ri % 2 == 0 else CW
        att, yds, td, fum = d['att'], d['yds'], d['td'], d['fum']
        sc(ws15, row, 1, player, bold=True, sz=10, fc=CW, bg=CR, h="left")
        sc(ws15, row, 2, att, sz=9, bg=bg, fmt="0")
        sc(ws15, row, 3, yds, bold=True, sz=9, fc="FF0E7060", bg="FFE8F8E8", fmt="0")
        sc(ws15, row, 4, round(yds / att, 1) if att else "", sz=9, bg=bg, fmt="0.0")
        sc(ws15, row, 5, d['lng'] if d['lng'] is not None else "", sz=9, bg=bg, fmt="0")
        sc(ws15, row, 6, _dash(td), sz=9, bg=bg, fmt="0" if td else "General")
        sc(ws15, row, 7, _dash(fum), sz=9, bg=bg, fmt="0" if fum else "General")
        row += 1
    if not rush_ranked:
        ws15.merge_cells(start_row=row, start_column=1, end_row=row, end_column=NC15)
        c = ws15.cell(row=row, column=1, value="No rusher data tagged (OPP RUSHER column).")
        c.font = Font(name=FN, sz=9, italic=True, color=CDG); c.alignment = Alignment(horizontal="center")
        row += 1
    row += 1

    banner(ws15, row, "RECEIVING", NC15, bg=CR, sz=13, ht=26); row += 1
    for c, txt in enumerate(["Player", "Rec", "Yds", "Avg", "Lng", "TD", "Fum", "Drop"], 1):
        hdr(ws15, row, c, txt, bg=CB, sz=9)
    row += 1
    rec_ranked = sorted(receiving.items(), key=lambda kv: -kv[1]['yds'])
    for ri, (player, d) in enumerate(rec_ranked):
        bg = CL if ri % 2 == 0 else CW
        rec, yds, td, fum, drop = d['rec'], d['yds'], d['td'], d['fum'], d['drop']
        sc(ws15, row, 1, player, bold=True, sz=10, fc=CW, bg=CBl, h="left")
        sc(ws15, row, 2, rec, sz=9, bg=bg, fmt="0")
        sc(ws15, row, 3, yds, bold=True, sz=9, fc="FF0E7060", bg="FFE8F8E8", fmt="0")
        sc(ws15, row, 4, round(yds / rec, 1) if rec else "", sz=9, bg=bg, fmt="0.0")
        sc(ws15, row, 5, d['lng'] if d['lng'] is not None else "", sz=9, bg=bg, fmt="0")
        sc(ws15, row, 6, _dash(td), sz=9, bg=bg, fmt="0" if td else "General")
        sc(ws15, row, 7, _dash(fum), sz=9, bg=bg, fmt="0" if fum else "General")
        sc(ws15, row, 8, _dash(drop), sz=9, bg=bg, fmt="0" if drop else "General")
        row += 1
    if not rec_ranked:
        ws15.merge_cells(start_row=row, start_column=1, end_row=row, end_column=NC15)
        c = ws15.cell(row=row, column=1, value="No receiver data tagged (OPP RECEIVER column).")
        c.font = Font(name=FN, sz=9, italic=True, color=CDG); c.alignment = Alignment(horizontal="center")
        row += 1

    ws15.freeze_panes = "A2"

    buf = io.BytesIO(); wb2.save(buf); buf.seek(0)
    return buf.getvalue()

# ── HTML Report ───────────────────────────────────────────────
def build_html(plays, opp, week, date):
    zc = {"BZ": "#c0392b", "OF": "#1a5276", "MF": "#0e7060", "FZ": "#7d6608", "RZ": "#c0392b", "GL": "#4a235a"}
    total = len(plays); runs = [p for p in plays if p['rp'] == 'Run']; passes = [p for p in plays if p['rp'] == 'Pass']

    def tags(items, cls=''):
        if not items: return '<span class="ctag">—</span>'
        return ''.join(f'<span class="ctag {cls}">{x["v"]} ({x["n"]})</span>' for x in items)

    zone_cards = ''
    for z in ZONE_LIST:
        zp = [p for p in plays if p['zone'] == z]
        if not zp: continue
        zr2 = [p for p in zp if p['rp'] == 'Run']; zpas = [p for p in zp if p['rp'] == 'Pass']
        rp = pct(len(zr2), len(zp)); pp = pct(len(zpas), len(zp))
        zone_cards += f'''<div class="zone-card">
          <div class="zone-hdr" style="background:{zc[z]}20;border-bottom:2px solid {zc[z]}">
            <div><div class="zone-badge" style="color:{zc[z]}">{z}</div><div class="zone-sub">{ZONE_NAMES[z]}</div></div>
            <div class="zone-plays">{len(zp)} plays</div>
          </div>
          <div class="zone-body">
            <div class="bar-row">
              <div class="bar-labels"><span style="color:#e8a095">RUN {rp}%</span><span style="color:#93d4f0">PASS {pp}%</span></div>
              <div class="bar-bg"><div class="bar-fill" style="background:#c0392b;width:{rp}%"></div></div>
            </div>
            <div class="zone-tags">
              <div class="tag-lbl">Top Run Concepts</div>{tags(top3(zr2,"concept"),"f")}
              <div class="tag-lbl" style="margin-top:6px">Top Pass Concepts</div>{tags(top3(zpas,"concept"),"c")}
              <div class="tag-lbl" style="margin-top:6px">Top Formations</div>{tags(top3(zp,"form"),"b")}
            </div>
          </div>
        </div>'''

    hash_cards = ''
    for h, lbl, cls, color in [('L', 'Left Hash', 'hl', '#b388d4'), ('M', 'Middle', 'hm', '#5dade2'), ('R', 'Right Hash', 'hr', '#e59866')]:
        hp = [p for p in plays if p['hash'] == h]
        if not hp:
            hash_cards += f'<div class="hash-card {cls}"><div class="hc-title" style="color:{color}">{lbl}</div><p style="color:rgba(240,237,232,.25);text-align:center;font-size:12px">No data</p></div>'
            continue
        hr2 = [p for p in hp if p['rp'] == 'Run']; hpass = [p for p in hp if p['rp'] == 'Pass']
        tf = top3(hp, 'concept', 1)
        hash_cards += f'''<div class="hash-card {cls}">
          <div class="hc-title" style="color:{color}">{lbl}</div>
          <div class="hbig" style="color:{color}">{len(hp)}</div>
          <div class="hsub">total plays</div>
          <div class="hrp">
            <div class="hrp-item"><div class="hrp-lbl">Run %</div><div class="hrp-val" style="color:#c0392b">{pct(len(hr2),len(hp))}%</div></div>
            <div class="hrp-item"><div class="hrp-lbl">Pass %</div><div class="hrp-val" style="color:#5dade2">{pct(len(hpass),len(hp))}%</div></div>
          </div>
          <div class="htc">Top Concept: <span style="color:#d4a017">{tf[0]["v"]+" ("+str(tf[0]["n"])+")" if tf else "—"}</span></div>
        </div>'''

    sit_rows = ''
    for lbl, fn in DD_SITS:
        sp = [p for p in plays if fn(p)]
        sr2 = [p for p in sp if p['rp'] == 'Run']; spass = [p for p in sp if p['rp'] == 'Pass']
        tf = top3(sr2, 'concept', 1); tc = top3(spass, 'concept', 1)
        sit_rows += f'''<tr>
          <td class="sit-lbl">{lbl}</td>
          <td style="text-align:center;font-family:Barlow Condensed,sans-serif;font-weight:800;font-size:20px;color:#c0392b">{pct(len(sr2),len(sp)) if sp else "—"}%</td>
          <td style="text-align:center;font-family:Barlow Condensed,sans-serif;font-weight:800;font-size:20px;color:#5dade2">{pct(len(spass),len(sp)) if sp else "—"}%</td>
          <td style="font-family:Share Tech Mono,monospace;font-size:9px;color:#d4a017">{tf[0]["v"]+" ("+str(tf[0]["n"])+")" if tf else "—"}</td>
          <td style="font-family:Share Tech Mono,monospace;font-size:9px;color:#d4a017">{tc[0]["v"]+" ("+str(tc[0]["n"])+")" if tc else "—"}</td>
        </tr>'''

    def con_rows(items, color, max_n):
        if not items: return '<div style="color:rgba(240,237,232,.3);font-size:12px">No data — tag concepts in Hudl to see trends</div>'
        return ''.join(f'<div style="padding:9px 0;border-bottom:1px solid rgba(240,237,232,.1)"><div style="display:flex;justify-content:space-between"><span style="font-family:Share Tech Mono,monospace;font-size:11px;color:{color}">{x["v"]}</span><span style="font-family:Barlow Condensed,sans-serif;font-weight:700;font-size:20px;color:{color}">{x["n"]}</span></div><div style="background:rgba(240,237,232,.06);height:3px;margin-top:4px"><div style="background:{color};height:3px;width:{round(x["n"]/max_n*100)}%"></div></div></div>' for x in items)

    mr = top3(runs, 'concept')[0]['n'] if top3(runs, 'concept') else 1
    mp = top3(passes, 'concept')[0]['n'] if top3(passes, 'concept') else 1
    mf = top3(plays, 'form')[0]['n'] if top3(plays, 'form') else 1

    biggest = compute_biggest_tendencies(plays, top_n=6)
    big_html = ''.join(f'<div style="padding:8px 0;border-bottom:1px solid rgba(240,237,232,.1);font-size:14px">{i+1}. {t}</div>' for i, t in enumerate(biggest)) or '<div style="color:rgba(240,237,232,.35);font-size:13px">Not enough tagged data yet.</div>'

    return f'''<!DOCTYPE html><html lang="en"><head><meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>DefensiveIQ — {opp}</title>
<link href="https://fonts.googleapis.com/css2?family=Oswald:wght@500;600;700;800&family=Inter:wght@400;500;600;700&display=swap" rel="stylesheet">
<style>:root{{--field:#0a1628;--chalk:#f0ede8;--red:#c0392b;--gold:#d4a017;--blue:#1a5276;--mid:#1e2d3d;--line:rgba(240,237,232,0.1);}}*{{box-sizing:border-box;margin:0;padding:0;}}body{{background:var(--field);color:var(--chalk);font-family:Inter,sans-serif;font-size:17px;line-height:1.5;}}nav{{display:flex;align-items:center;justify-content:space-between;padding:14px 40px;border-bottom:1px solid var(--line);background:rgba(10,22,40,.97);}}.logo{{font-family:Oswald,sans-serif;font-weight:900;font-size:22px;}}.logo span{{color:#c0392b;}}.wrap{{max-width:1200px;margin:0 auto;padding:40px;}}.eyebrow{{font-family:Inter,sans-serif;font-size:19px;letter-spacing:.2em;color:var(--gold);text-transform:uppercase;margin-bottom:10px;}}.rpt-hdr{{display:flex;justify-content:space-between;align-items:flex-end;margin-bottom:32px;padding-bottom:18px;border-bottom:1px solid var(--line);}}.rpt-title{{font-family:Oswald,sans-serif;font-weight:900;font-size:42px;text-transform:uppercase;}}.rpt-meta{{font-family:Inter,sans-serif;font-size:19px;color:var(--gold);text-align:right;}}.sum-grid{{display:grid;grid-template-columns:repeat(5,1fr);gap:10px;margin-bottom:36px;}}.sum-card{{background:var(--mid);border:1px solid var(--line);padding:16px;}}.sum-lbl{{font-size:13px;font-weight:600;letter-spacing:.14em;text-transform:uppercase;color:rgba(240,237,232,.65);margin-bottom:5px;}}.sum-val{{font-family:Oswald,sans-serif;font-weight:800;font-size:37px;line-height:1;}}.stitle{{font-family:Oswald,sans-serif;font-weight:800;font-size:21px;text-transform:uppercase;letter-spacing:.06em;margin-bottom:16px;margin-top:36px;padding-bottom:8px;border-bottom:1px solid var(--line);}}.zone-grid{{display:grid;grid-template-columns:repeat(3,1fr);gap:12px;}}.zone-card{{background:var(--mid);border:1px solid var(--line);overflow:hidden;}}.zone-hdr{{padding:10px 14px;display:flex;justify-content:space-between;align-items:center;}}.zone-badge{{font-family:Oswald,sans-serif;font-weight:900;font-size:19px;}}.zone-sub{{font-size:13px;color:rgba(240,237,232,.65);}}.zone-plays{{font-family:Inter,sans-serif;font-size:13px;color:var(--gold);}}.zone-body{{padding:11px 14px;}}.bar-row{{margin-bottom:8px;}}.bar-labels{{display:flex;justify-content:space-between;font-size:13px;font-weight:600;letter-spacing:.06em;text-transform:uppercase;margin-bottom:3px;}}.bar-bg{{background:rgba(240,237,232,.07);height:5px;position:relative;}}.bar-fill{{height:5px;position:absolute;left:0;top:0;}}.zone-tags{{margin-top:8px;border-top:1px solid var(--line);padding-top:8px;}}.tag-lbl{{font-size:12px;font-weight:600;letter-spacing:.1em;text-transform:uppercase;color:rgba(240,237,232,.82);margin-bottom:3px;}}.ctag{{display:inline-block;background:rgba(240,237,232,.05);border:1px solid rgba(240,237,232,.1);font-family:Inter,sans-serif;font-size:12px;padding:2px 4px;margin:1px 1px 1px 0;color:rgba(240,237,232,.82);}}.ctag.f{{border-color:rgba(192,57,43,.4);color:#e8a095;}}.ctag.c{{border-color:rgba(93,173,226,.35);color:#93d4f0;}}.ctag.b{{border-color:rgba(212,160,23,.4);color:#d4a017;}}.hash-grid{{display:grid;grid-template-columns:repeat(3,1fr);gap:12px;}}.hash-card{{background:var(--mid);border:1px solid var(--line);padding:18px;text-align:center;}}.hc-title{{font-family:Oswald,sans-serif;font-weight:800;font-size:15px;letter-spacing:.14em;text-transform:uppercase;margin-bottom:12px;padding-bottom:8px;border-bottom:1px solid var(--line);}}.hbig{{font-family:Oswald,sans-serif;font-weight:900;font-size:50px;line-height:1;margin-bottom:2px;}}.hsub{{font-size:13px;color:rgba(240,237,232,.62);margin-bottom:10px;}}.hrp{{display:grid;grid-template-columns:1fr 1fr;gap:6px;margin-bottom:8px;}}.hrp-item{{background:rgba(240,237,232,.04);padding:7px;}}.hrp-lbl{{font-size:12px;letter-spacing:.08em;text-transform:uppercase;color:rgba(240,237,232,.62);}}.hrp-val{{font-family:Oswald,sans-serif;font-weight:700;font-size:21px;}}.htc{{font-family:Inter,sans-serif;font-size:13px;color:rgba(240,237,232,.62);}}.sit-table{{width:100%;border-collapse:collapse;font-size:15px;}}.sit-table th{{background:var(--field);padding:7px 10px;font-family:Oswald,sans-serif;font-weight:700;font-size:13px;letter-spacing:.1em;text-transform:uppercase;color:rgba(240,237,232,.7);border:1px solid var(--line);text-align:center;}}.sit-table th:first-child{{text-align:left;}}.sit-table td{{border:1px solid var(--line);padding:8px 12px;}}.sit-table tr:nth-child(odd) td{{background:rgba(240,237,232,.02);}}.sit-table tr:nth-child(even) td{{background:var(--mid);}}.sit-lbl{{font-family:Oswald,sans-serif;font-weight:700;font-size:15px;white-space:nowrap;}}.con-grid{{display:grid;grid-template-columns:repeat(3,1fr);gap:24px;}}.alert-box{{background:var(--mid);border:1px solid var(--line);border-left:4px solid var(--red);padding:16px 20px;}}</style>
</head><body>
<nav><div class="logo">DEFENSIVE<span>IQ</span></div><div style="font-family:Share Tech Mono,monospace;font-size:10px;color:rgba(240,237,232,.4)">OPPONENT OFFENSIVE TENDENCY REPORT</div></nav>
<div class="wrap">
<div class="rpt-hdr"><div><div class="eyebrow">// Defensive Coordinator — Opponent Scouting Report</div><div class="rpt-title">{opp} — Offensive Analysis</div></div><div class="rpt-meta">WEEK {week}{("<br>"+date) if date else ""}<br>{total} PLAYS ANALYZED</div></div>
<div class="sum-grid">
  <div class="sum-card"><div class="sum-lbl">Total Plays</div><div class="sum-val">{total}</div></div>
  <div class="sum-card"><div class="sum-lbl">Run %</div><div class="sum-val" style="color:#c0392b">{pct(len(runs),total)}%</div></div>
  <div class="sum-card"><div class="sum-lbl">Pass %</div><div class="sum-val" style="color:#5dade2">{pct(len(passes),total)}%</div></div>
  <div class="sum-card"><div class="sum-lbl">Avg Yds/Play</div><div class="sum-val" style="color:#d4a017">{_avg(plays):.1f}</div></div>
  <div class="sum-card"><div class="sum-lbl">Expl. Play %</div><div class="sum-val" style="color:#0e7060">{pct(len([p for p in plays if p["expl"]]),total)}%</div></div>
</div>
<div class="stitle">Biggest Tendencies</div>
<div class="alert-box">{big_html}</div>
<div class="stitle">Field Zone Breakdown</div><div class="zone-grid">{zone_cards}</div>
<div class="stitle">Hash Tendencies</div><div class="hash-grid">{hash_cards}</div>
<div class="stitle">Down &amp; Distance Summary</div>
<table class="sit-table">
  <tr><th style="text-align:left">Situation</th><th>Run%</th><th>Pass%</th><th style="text-align:left">Top Run</th><th style="text-align:left">Top Pass</th></tr>
  {sit_rows}
</table>
<div class="stitle">Concept Frequency</div>
<div class="con-grid">
  <div><div class="eyebrow" style="margin-bottom:12px">// Top Run Concepts</div>{con_rows(top3(runs,"concept"),"#e8a095",mr)}</div>
  <div><div class="eyebrow" style="margin-bottom:12px">// Top Pass Concepts</div>{con_rows(top3(passes,"concept"),"#93d4f0",mp)}</div>
  <div><div class="eyebrow" style="margin-bottom:12px">// Top Formations</div>{con_rows(top3(plays,"form"),"#d4a017",mf)}</div>
</div>
</div></body></html>'''

# ── STREAMLIT UI ──────────────────────────────────────────────
st.markdown('<div class="main-title">Defensive<span style="color:#c0392b">IQ</span></div>', unsafe_allow_html=True)
st.markdown('<div style="font-size:16px;color:rgba(240,237,232,.55);margin-bottom:24px;font-weight:300">Scout your next opponent\'s offense. Upload their playlist and uncover every tendency — formations, concepts, runs, passes, hashes, and situations — to build your defensive game plan.</div>', unsafe_allow_html=True)

st.divider()
col1, col2, col3 = st.columns(3)
with col1: opp = st.text_input("Opponent Name", placeholder="e.g. Lincoln High School")
with col2: week = st.text_input("Week", placeholder="e.g. 11")
with col3: date = st.text_input("Game Date", placeholder="e.g. Nov 7, 2026")

st.markdown("**Presentation Colors** — used for your scouting PowerPoint")
cc1, cc2 = st.columns(2)
with cc1: team_primary = st.color_picker("Primary (headers/titles)", "#16213E")
with cc2: team_accent = st.color_picker("Accent (highlights)", "#C0392B")

st.markdown("---")
uploaded = st.file_uploader("Upload Opponent Offensive Playlist (.xlsx or .csv)", type=['xlsx', 'xls', 'csv'],
                             help="Export the opponent's offensive playlist from Hudl as Excel or CSV and upload here.")

if uploaded and st.button("🛡️ RUN ANALYSIS"):
    with st.spinner("Analyzing opponent offensive tendencies..."):
        try:
            if uploaded.name.lower().endswith('.csv'):
                df = pd.read_csv(uploaded)
            else:
                df = pd.read_excel(uploaded)

            df, matched, missing = map_columns(df)
            with st.expander("📋 Column mapping — what we found in your file"):
                st.write("**Matched:** " + (", ".join(matched.keys()) if matched else "none"))
                if missing:
                    st.info("Not found (those sections will be blank): " + ", ".join(missing))

            req_missing = check_required(matched)
            if req_missing:
                st.error("Your file is missing required column(s): " + ", ".join(req_missing))
                st.info("These are needed to analyze anything. Check that your export includes "
                        "Play Type (Run/Pass), Yard Line, Down, and Distance — then re-export and try again.")
                st.stop()

            plays = load_plays(df)
            if len(plays) == 0:
                st.error("No Run/Pass plays found in this file.")
                st.info("Common causes: PLAY TYPE uses different words than 'Run'/'Pass', or YARD LN is blank. "
                        "Open the column mapping above to see what we detected.")
            else:
                opp_name = opp or "Opponent"
                prog = st.progress(0, "Reading data...")
                prog.progress(30, "Calculating field-zone tendencies...")
                excel_bytes = build_excel(plays, opp_name, week, date)
                prog.progress(65, "Building HTML report...")
                html_bytes = build_html(plays, opp_name, week, date).encode('utf-8')
                prog.progress(85, "Building scouting presentation...")
                pptx_bytes = build_pptx(plays, opp_name, week, date, team_primary, team_accent)
                prog.progress(100, "Complete!")

                runs = [p for p in plays if p['rp'] == 'Run']; passes = [p for p in plays if p['rp'] == 'Pass']

                st.success(f"✅ Analysis complete — {len(plays)} plays analyzed")
                st.divider()

                m1, m2, m3, m4, m5 = st.columns(5)
                m1.metric("Total Plays", len(plays))
                m2.metric("Run %", f"{pct(len(runs), len(plays))}%")
                m3.metric("Pass %", f"{pct(len(passes), len(plays))}%")
                m4.metric("Avg Yds/Play", f"{_avg(plays):.1f}")
                m5.metric("Explosive %", f"{pct(len([p for p in plays if p['expl']]), len(plays))}%")

                st.divider()
                st.markdown("### Download Your Reports")
                d1, d2, d3 = st.columns(3)
                fname = (opp_name + "_" if opp_name else "") + (f"Week{week}_" if week else "") + "DefensiveIQ"
                with d1:
                    st.download_button("📊 Excel Workbook", data=excel_bytes,
                        file_name=f"{fname}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
                with d2:
                    st.download_button("🌐 HTML Report", data=html_bytes,
                        file_name=f"{fname}_Report.html", mime="text/html")
                with d3:
                    st.download_button("📽️ Scouting Presentation", data=pptx_bytes,
                        file_name=f"{fname}_Scouting.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

                st.divider()
                st.markdown("### Quick Summary")
                z1, z2, z3 = st.columns(3)
                with z1:
                    st.markdown("**Top Run Concepts**")
                    for x in top3(runs, 'concept'): st.markdown(f"- {x['v']} ({x['n']} calls)")
                    if not runs: st.markdown("*No run plays tagged*")
                with z2:
                    st.markdown("**Top Pass Concepts**")
                    for x in top3(passes, 'concept'): st.markdown(f"- {x['v']} ({x['n']} calls)")
                    if not passes: st.markdown("*No pass plays tagged*")
                with z3:
                    st.markdown("**Top Formations**")
                    for x in top3(plays, 'form'): st.markdown(f"- {x['v']} ({x['n']} snaps)")
                    if not any(p['form'] for p in plays): st.markdown("*Tag OFF FORM in Hudl*")

                st.markdown("### Biggest Tendencies")
                big = compute_biggest_tendencies(plays, top_n=6)
                if big:
                    for i, t in enumerate(big, 1): st.markdown(f"**#{i}** — {t}")
                else:
                    st.markdown("*Not enough tagged snaps yet to compute headline tendencies.*")

        except Exception as e:
            import traceback
            st.error("Something went wrong reading this file — it may be formatted differently than expected.")
            st.info("Check that this is a Hudl playlist export with PLAY TYPE, YARD LN, OFF FORM, DN, DIST, "
                    "HASH, and a concept column (OFF PLAY / PASS CONCEPT).")
            with st.expander("🔧 Technical details"):
                st.code(traceback.format_exc())

st.divider()
st.markdown('<div style="font-family:Share Tech Mono,monospace;font-size:10px;color:rgba(240,237,232,.25);text-align:center;padding:20px 0">© 2026 DEFENSIVEIQ · BUILT FOR DEFENSIVE COORDINATORS</div>', unsafe_allow_html=True)